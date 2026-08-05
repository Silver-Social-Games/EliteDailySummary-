"""Shared VIP Event PowerPoint theme — Vegas 2026 gold/dark styling."""
from __future__ import annotations

import re

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

GOLD = "FFCD36"
BLACK = "000000"
DARK = "1A1A1A"
DARK_CARD = "252525"
DARK_ALT = "2E2E2E"
DARK_VALUE = "2A2A2A"
GOLD_MUTED = "3D3520"
INK = "1A1A1A"
MUTED = "AAAAAA"
WHITE = "FFFFFF"
BEIGE = "F5F0E6"
SURFACE_ALT = "F5F5F5"
CHART_GRAY = "666666"
BODY_TEXT = WHITE

# Legacy aliases kept for imports that still reference them
NAVY = DARK
ORANGE = GOLD

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TITLE_MAIN = "VIP Event · Vegas 2026"

# Safe content area (inches) — keeps everything inside slide when copied
MARGIN_X = 0.36
MARGIN_BOTTOM = 0.28
SLIDE_W_IN = float(SLIDE_W / Inches(1))
SLIDE_H_IN = float(SLIDE_H / Inches(1))
CONTENT_RIGHT = SLIDE_W_IN - MARGIN_X
CONTENT_BOTTOM = SLIDE_H_IN - MARGIN_BOTTOM

# Soft pastel chart palette
CHART_COLORS = (
    "A8D5BA",
    "F4C2C2",
    "B8C9E8",
    "E8D4A8",
    "D4B8E8",
    "C2E8F4",
    "E8C2D4",
    "B8E8C2",
    "E8B8B8",
)


def rgb(hex6: str) -> RGBColor:
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def load_template(path: str | None) -> Presentation:
    prs = Presentation(path) if path else Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def set_slide_dark_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(DARK)


def add_title_band(slide, title: str = TITLE_MAIN, subtitle: str = "") -> None:
    set_slide_dark_background(slide)
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.35))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(DARK)
    band.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(11), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = rgb(GOLD)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.86), Inches(11), Inches(0.42))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.bold = True
        sp.font.color.rgb = rgb(GOLD)


def style_table_header(cell, text: str) -> None:
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb(GOLD)
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(8)
        p.font.color.rgb = rgb(BLACK)
        p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def style_table_body(
    cell,
    text: str,
    *,
    bold: bool = False,
    size: int = 8,
    alt_row: bool = False,
) -> None:
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb(DARK_ALT if alt_row else DARK_CARD)
    for p in cell.text_frame.paragraphs:
        p.font.bold = bold
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(BODY_TEXT)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_speaker_notes(slide, items: list[tuple[str, str]]) -> None:
    """Populate presenter notes from (slide_line, speaker_note) pairs."""
    if not items:
        return
    tf = slide.notes_slide.notes_text_frame
    for i, (slide_line, note) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i + 1}. {slide_line}\n{note}"
        p.font.size = Pt(11)


def add_executive_bullets(
    slide,
    items: list[tuple[str, str]],
    *,
    top: float = 1.52,
) -> None:
    """C-level bullet list — short on-slide lines only; pair with set_speaker_notes."""
    n = len(items)
    if n == 0:
        return
    left = MARGIN_X
    width = CONTENT_RIGHT - MARGIN_X
    usable_h = CONTENT_BOTTOM - top - 0.15
    row_h = min(0.72, usable_h / n)
    num_col = 0.42
    text_left = left + num_col + 0.12
    text_w = width - num_col - 0.12
    line_pt = 15 if n <= 5 else 13

    for i, (line, _) in enumerate(items):
        y = top + i * row_h + row_h * 0.12

        nbox = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(num_col), Inches(row_h * 0.7))
        np = nbox.text_frame.paragraphs[0]
        np.text = str(i + 1)
        np.font.size = Pt(18 if n <= 5 else 16)
        np.font.bold = True
        np.font.color.rgb = rgb(GOLD)
        np.alignment = PP_ALIGN.RIGHT

        tbox = slide.shapes.add_textbox(Inches(text_left), Inches(y), Inches(text_w), Inches(row_h * 0.85))
        ttf = tbox.text_frame
        ttf.word_wrap = True
        ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if " — " in line:
            head, tail = line.split(" — ", 1)
            tp = ttf.paragraphs[0]
            tp.text = head
            tp.font.size = Pt(line_pt)
            tp.font.bold = True
            tp.font.color.rgb = rgb(GOLD)
            if tail:
                sp = ttf.add_paragraph()
                sp.text = tail
                sp.font.size = Pt(line_pt - 2)
                sp.font.bold = False
                sp.font.color.rgb = rgb(BODY_TEXT)
        else:
            tp = ttf.paragraphs[0]
            tp.text = line
            tp.font.size = Pt(line_pt)
            tp.font.bold = True
            tp.font.color.rgb = rgb(BODY_TEXT)


def _render_executive_bullet_lines(
    slide,
    items: list[tuple[str, str]],
    *,
    top: float,
    start_num: int = 1,
    line_pt: int = 13,
) -> float:
    """Render numbered executive bullets; returns bottom y (inches)."""
    n = len(items)
    if n == 0:
        return top
    left = MARGIN_X
    width = CONTENT_RIGHT - MARGIN_X
    row_h = min(0.58, (CONTENT_BOTTOM - top - 0.1) / max(n, 1))
    num_col = 0.38
    text_left = left + num_col + 0.1
    text_w = width - num_col - 0.1

    for i, (line, _) in enumerate(items):
        y = top + i * row_h + row_h * 0.08
        nbox = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(num_col), Inches(row_h * 0.75))
        np = nbox.text_frame.paragraphs[0]
        np.text = str(start_num + i)
        np.font.size = Pt(15 if n <= 4 else 14)
        np.font.bold = True
        np.font.color.rgb = rgb(GOLD)
        np.alignment = PP_ALIGN.RIGHT

        tbox = slide.shapes.add_textbox(Inches(text_left), Inches(y), Inches(text_w), Inches(row_h * 0.9))
        ttf = tbox.text_frame
        ttf.word_wrap = True
        ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if " — " in line:
            head, tail = line.split(" — ", 1)
            tp = ttf.paragraphs[0]
            tp.text = head
            tp.font.size = Pt(line_pt)
            tp.font.bold = True
            tp.font.color.rgb = rgb(GOLD)
            if tail:
                sp = ttf.add_paragraph()
                sp.text = tail
                sp.font.size = Pt(line_pt - 2)
                sp.font.bold = False
                sp.font.color.rgb = rgb(BODY_TEXT)
        else:
            tp = ttf.paragraphs[0]
            tp.text = line
            tp.font.size = Pt(line_pt)
            tp.font.bold = True
            tp.font.color.rgb = rgb(BODY_TEXT)

    return top + n * row_h


def add_section_heading(slide, text: str, top: float) -> float:
    """Gold section label; returns y below heading."""
    box = slide.shapes.add_textbox(
        Inches(MARGIN_X), Inches(top), Inches(CONTENT_RIGHT - MARGIN_X), Inches(0.28)
    )
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = rgb(GOLD)
    return top + 0.32


def _goal_line_text(item: str | tuple[str, str]) -> str:
    return item[0] if isinstance(item, tuple) else str(item)


def _is_placeholder_line(text: str) -> bool:
    return text.lower().startswith("main goal") or text.lower().startswith("secondary goal")


def _label_in_box(slide, left, top, width, height, text: str, *, hero: bool = False, banner: bool = False) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    if banner:
        box.fill.fore_color.rgb = rgb(GOLD_MUTED)
    elif hero:
        box.fill.fore_color.rgb = rgb(DARK_VALUE)
    else:
        box.fill.fore_color.rgb = rgb(DARK_CARD)
    box.line.color.rgb = rgb(GOLD if hero or banner else GOLD_MUTED)
    box.line.width = Pt(1.5 if hero else 0.9)
    if box.adjustments:
        box.adjustments[0] = 0.06

    tbox = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(top + 0.08), Inches(width - 0.24), Inches(height - 0.16)
    )
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14 if hero else (11 if banner else 10))
    p.font.bold = hero or banner
    p.font.italic = _is_placeholder_line(text)
    p.font.color.rgb = rgb(MUTED if _is_placeholder_line(text) else (GOLD if hero else BODY_TEXT))
    p.alignment = PP_ALIGN.CENTER


GOAL_MAIN_OUTLINES = (GOLD, "A8D5BA", "B8C9E8")
GOAL_SEC_OUTLINES = ("E8D4A8", "F4C2C2", "D4B8E8")


def _outline_shape(shape, color_hex: str, *, pt: float = 2.5) -> None:
    shape.line.color.rgb = rgb(color_hex)
    shape.line.width = Pt(pt)


def _secondary_cluster_x(count: int, index: int, item_w: float, gap: float, center_x: float) -> float:
    """X position (left edge) for item in a centered tight row."""
    total = count * item_w + (count - 1) * gap
    start = center_x - total / 2
    return start + index * (item_w + gap)


def _goal_text_in_shape(slide, left, top, width, height, text: str, *, size: int = 10, gold: bool = False) -> None:
    tbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = gold
    p.font.italic = _is_placeholder_line(text)
    p.font.color.rgb = rgb(MUTED if _is_placeholder_line(text) else (GOLD if gold else BODY_TEXT))
    p.alignment = PP_ALIGN.CENTER


def add_goals_layout_option_a(
    slide,
    main_lines: list[str],
    secondary_lines: list[str],
    *,
    top: float = 1.48,
) -> None:
    """Option A — Slot reels: three equal reel columns + coin-stack secondary row."""
    usable_w = CONTENT_RIGHT - MARGIN_X
    cx = MARGIN_X + usable_w / 2
    base_y = 4.92

    cap = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(top), Inches(4.5), Inches(0.22))
    cap.text_frame.paragraphs[0].text = "MAIN GOALS · THE REELS"
    cap.text_frame.paragraphs[0].font.size = Pt(9)
    cap.text_frame.paragraphs[0].font.bold = True
    cap.text_frame.paragraphs[0].font.color.rgb = rgb(GOLD)

    frame_left = MARGIN_X + 0.06
    frame_w = usable_w - 0.12
    frame_top = top + 0.34
    frame_h = base_y - frame_top
    frame_pad = 0.16

    cabinet = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(frame_left),
        Inches(frame_top),
        Inches(frame_w),
        Inches(frame_h),
    )
    cabinet.fill.solid()
    cabinet.fill.fore_color.rgb = rgb(DARK_ALT)
    _outline_shape(cabinet, GOLD, pt=3.0)
    if cabinet.adjustments:
        cabinet.adjustments[0] = 0.04

    marquee = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(frame_left + frame_pad),
        Inches(frame_top + 0.08),
        Inches(frame_w - 2 * frame_pad),
        Inches(0.28),
    )
    marquee.fill.solid()
    marquee.fill.fore_color.rgb = rgb(GOLD)
    marquee.line.fill.background()
    if marquee.adjustments:
        marquee.adjustments[0] = 0.2

    marquee_txt = slide.shapes.add_textbox(
        Inches(frame_left + frame_pad),
        Inches(frame_top + 0.1),
        Inches(frame_w - 2 * frame_pad),
        Inches(0.24),
    )
    marquee_txt.text_frame.paragraphs[0].text = "★  SPIN TO WIN  ★"
    marquee_txt.text_frame.paragraphs[0].font.size = Pt(9)
    marquee_txt.text_frame.paragraphs[0].font.bold = True
    marquee_txt.text_frame.paragraphs[0].font.color.rgb = rgb(BLACK)
    marquee_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    inner_left = frame_left + frame_pad
    inner_w = frame_w - 2 * frame_pad
    inner_top = frame_top + frame_pad + 0.28
    inner_h = frame_h - 2 * frame_pad - 0.28
    reel_gap = 0.12
    reel_w = (inner_w - 2 * reel_gap) / 3
    reel_h = inner_h

    for i in range(3):
        line = main_lines[i] if i < len(main_lines) else f"Main goal {i + 1}"
        outline = GOAL_MAIN_OUTLINES[i % len(GOAL_MAIN_OUTLINES)]
        x = inner_left + i * (reel_w + reel_gap)
        y = inner_top

        glow = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x - 0.025),
            Inches(y - 0.025),
            Inches(reel_w + 0.05),
            Inches(reel_h + 0.05),
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = rgb(outline)
        glow.line.fill.background()
        if glow.adjustments:
            glow.adjustments[0] = 0.06

        body = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(reel_w), Inches(reel_h)
        )
        body.fill.solid()
        body.fill.fore_color.rgb = rgb(DARK_CARD)
        _outline_shape(body, outline, pt=2.5)
        if body.adjustments:
            body.adjustments[0] = 0.08

        badge_d = 0.32
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + reel_w / 2 - badge_d / 2),
            Inches(y + 0.06),
            Inches(badge_d),
            Inches(badge_d),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(GOLD)
        badge.line.color.rgb = rgb(outline)
        badge.line.width = Pt(1.5)

        num = slide.shapes.add_textbox(
            Inches(x + reel_w / 2 - badge_d / 2),
            Inches(y + 0.1),
            Inches(badge_d),
            Inches(0.24),
        )
        num.text_frame.paragraphs[0].text = str(i + 1)
        num.text_frame.paragraphs[0].font.size = Pt(12)
        num.text_frame.paragraphs[0].font.bold = True
        num.text_frame.paragraphs[0].font.color.rgb = rgb(BLACK)
        num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        sym_top = y + 0.48
        sym_h = reel_h * 0.38
        sym_w = reel_w * 0.2
        sym_gap = (reel_w - 3 * sym_w) / 4
        sym_colors = (outline, GOLD_MUTED, outline)
        for j in range(3):
            sx = x + sym_gap + j * (sym_w + sym_gap)
            sym = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(sx),
                Inches(sym_top),
                Inches(sym_w),
                Inches(sym_h),
            )
            sym.fill.solid()
            sym.fill.fore_color.rgb = rgb(sym_colors[j % len(sym_colors)])
            sym.line.color.rgb = rgb(GOLD_MUTED)
            sym.line.width = Pt(0.75)
            if sym.adjustments:
                sym.adjustments[0] = 0.15

        divider_y = sym_top + sym_h + 0.06
        divider = slide.shapes.add_shape(
            1, Inches(x + 0.06), Inches(divider_y), Inches(reel_w - 0.12), Inches(0.03)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = rgb(GOLD)
        divider.line.fill.background()

        win_top = divider_y + 0.1
        win_h = y + reel_h - win_top - 0.1
        window = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.08),
            Inches(win_top),
            Inches(reel_w - 0.16),
            Inches(win_h),
        )
        window.fill.solid()
        window.fill.fore_color.rgb = rgb(DARK_VALUE)
        window.line.color.rgb = rgb(GOLD_MUTED)
        window.line.width = Pt(1.0)
        if window.adjustments:
            window.adjustments[0] = 0.1

        _goal_text_in_shape(
            slide, x + 0.1, win_top + 0.06, reel_w - 0.2, win_h - 0.12, line, size=9
        )

    sec_cap = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(base_y + 0.1), Inches(4.5), Inches(0.2))
    sec_cap.text_frame.paragraphs[0].text = "SECONDARY GOALS · COIN STACK"
    sec_cap.text_frame.paragraphs[0].font.size = Pt(8)
    sec_cap.text_frame.paragraphs[0].font.color.rgb = rgb(MUTED)

    stack_w = 1.12
    sec_gap = 0.08
    sec_y = base_y + 0.38
    coin_d = 0.38
    for i in range(3):
        sec_line = secondary_lines[i] if i < len(secondary_lines) else f"Secondary goal {i + 1}"
        outline = GOAL_SEC_OUTLINES[i % len(GOAL_SEC_OUTLINES)]
        x = _secondary_cluster_x(3, i, stack_w, sec_gap, cx)
        base_cy = sec_y + 0.72

        for k in range(2):
            offset = k * 0.14
            coin = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + stack_w / 2 - coin_d / 2),
                Inches(base_cy - coin_d - offset),
                Inches(coin_d),
                Inches(coin_d * 0.22),
            )
            coin.fill.solid()
            coin.fill.fore_color.rgb = rgb(outline if k == 0 else GOLD_MUTED)
            coin.line.color.rgb = rgb(GOLD)
            coin.line.width = Pt(1.0)

        base_coin = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + 0.06),
            Inches(sec_y),
            Inches(stack_w - 0.12),
            Inches(stack_w - 0.12),
        )
        base_coin.fill.solid()
        base_coin.fill.fore_color.rgb = rgb(DARK_VALUE)
        _outline_shape(base_coin, outline, pt=2.0)
        _goal_text_in_shape(
            slide, x + 0.08, sec_y + 0.18, stack_w - 0.16, stack_w - 0.36, sec_line, size=8
        )


def add_goals_layout_option_b(
    slide,
    main_lines: list[str],
    secondary_lines: list[str],
    *,
    top: float = 1.5,
) -> None:
    """Option B — Jackpot hub: three equal jackpot meters + bet-line chip row."""
    usable_w = CONTENT_RIGHT - MARGIN_X
    cx = MARGIN_X + usable_w / 2

    cap = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(top), Inches(4.5), Inches(0.22))
    cap.text_frame.paragraphs[0].text = "MAIN GOALS · JACKPOT PANEL"
    cap.text_frame.paragraphs[0].font.size = Pt(9)
    cap.text_frame.paragraphs[0].font.bold = True
    cap.text_frame.paragraphs[0].font.color.rgb = rgb(GOLD)

    panel_gap = 0.16
    panel_w = (usable_w - 0.24 - 2 * panel_gap) / 3
    panel_h = 2.55
    panel_y = top + 0.38
    row_w = 3 * panel_w + 2 * panel_gap
    row_left = cx - row_w / 2

    hub = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(row_left - 0.12),
        Inches(panel_y - 0.1),
        Inches(row_w + 0.24),
        Inches(panel_h + 0.2),
    )
    hub.fill.solid()
    hub.fill.fore_color.rgb = rgb(DARK_ALT)
    _outline_shape(hub, GOLD, pt=2.5)
    if hub.adjustments:
        hub.adjustments[0] = 0.03

    for i in range(3):
        line = main_lines[i] if i < len(main_lines) else f"Main goal {i + 1}"
        outline = GOAL_MAIN_OUTLINES[i % len(GOAL_MAIN_OUTLINES)]
        x = row_left + i * (panel_w + panel_gap)
        y = panel_y

        glow = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x - 0.025),
            Inches(y - 0.025),
            Inches(panel_w + 0.05),
            Inches(panel_h + 0.05),
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = rgb(outline)
        glow.line.fill.background()
        if glow.adjustments:
            glow.adjustments[0] = 0.05

        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(panel_w), Inches(panel_h)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = rgb(DARK_CARD)
        _outline_shape(panel, outline, pt=2.5)
        if panel.adjustments:
            panel.adjustments[0] = 0.06

        hdr_h = 0.36
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(panel_w), Inches(hdr_h)
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = rgb(GOLD)
        hdr.line.fill.background()
        if hdr.adjustments:
            hdr.adjustments[0] = 0.15

        hdr_txt = slide.shapes.add_textbox(Inches(x), Inches(y + 0.06), Inches(panel_w), Inches(0.24))
        hdr_txt.text_frame.paragraphs[0].text = f"JACKPOT {i + 1}"
        hdr_txt.text_frame.paragraphs[0].font.size = Pt(9)
        hdr_txt.text_frame.paragraphs[0].font.bold = True
        hdr_txt.text_frame.paragraphs[0].font.color.rgb = rgb(BLACK)
        hdr_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        meter_top = y + hdr_h + 0.14
        meter_w = panel_w - 0.24
        meter_x = x + 0.12
        bar_h = 0.14
        bar_gap = 0.08
        fill_levels = (0.85, 0.65, 0.45)
        for j, fill_pct in enumerate(fill_levels):
            by = meter_top + j * (bar_h + bar_gap)
            track = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(meter_x),
                Inches(by),
                Inches(meter_w),
                Inches(bar_h),
            )
            track.fill.solid()
            track.fill.fore_color.rgb = rgb(DARK_VALUE)
            track.line.color.rgb = rgb(GOLD_MUTED)
            track.line.width = Pt(0.75)
            if track.adjustments:
                track.adjustments[0] = 0.3

            fill_w = meter_w * fill_pct
            fill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(meter_x),
                Inches(by),
                Inches(fill_w),
                Inches(bar_h),
            )
            fill.fill.solid()
            fill.fill.fore_color.rgb = rgb(outline if j == 0 else GOLD_MUTED)
            fill.line.fill.background()
            if fill.adjustments:
                fill.adjustments[0] = 0.3

        disp_top = meter_top + 3 * (bar_h + bar_gap) + 0.1
        disp_h = y + panel_h - disp_top - 0.12
        display = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(meter_x),
            Inches(disp_top),
            Inches(meter_w),
            Inches(disp_h),
        )
        display.fill.solid()
        display.fill.fore_color.rgb = rgb(DARK_VALUE)
        display.line.color.rgb = rgb(outline)
        display.line.width = Pt(1.5)
        if display.adjustments:
            display.adjustments[0] = 0.08

        _goal_text_in_shape(
            slide, meter_x + 0.04, disp_top + 0.06, meter_w - 0.08, disp_h - 0.12, line, size=9, gold=True
        )

    sec_y = 4.95
    sec_cap = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(sec_y - 0.28), Inches(5), Inches(0.2))
    sec_cap.text_frame.paragraphs[0].text = "SECONDARY GOALS · BET LINE"
    sec_cap.text_frame.paragraphs[0].font.size = Pt(8)
    sec_cap.text_frame.paragraphs[0].font.color.rgb = rgb(MUTED)

    chip_w, chip_h = 1.08, 0.68
    chip_gap = 0.1
    chip_row_y = sec_y + 0.05
    line_y = sec_y - 0.06
    payline = slide.shapes.add_shape(
        1,
        Inches(cx - 1.75),
        Inches(line_y),
        Inches(3.5),
        Inches(0.03),
    )
    payline.fill.solid()
    payline.fill.fore_color.rgb = rgb(GOLD)
    payline.line.fill.background()

    for i in range(3):
        line = secondary_lines[i] if i < len(secondary_lines) else f"Secondary goal {i + 1}"
        outline = GOAL_SEC_OUTLINES[i % len(GOAL_SEC_OUTLINES)]
        px = _secondary_cluster_x(3, i, chip_w, chip_gap, cx)

        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(px),
            Inches(chip_row_y),
            Inches(chip_w),
            Inches(chip_h),
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = rgb(DARK_CARD)
        _outline_shape(chip, outline, pt=2.0)
        if chip.adjustments:
            chip.adjustments[0] = 0.25

        inner = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(px + 0.06),
            Inches(chip_row_y + 0.06),
            Inches(chip_w - 0.12),
            Inches(chip_h - 0.12),
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = rgb(DARK_VALUE)
        inner.line.color.rgb = rgb(outline)
        inner.line.width = Pt(1.0)
        if inner.adjustments:
            inner.adjustments[0] = 0.3

        _goal_text_in_shape(
            slide, px + 0.06, chip_row_y + 0.12, chip_w - 0.12, chip_h - 0.22, line, size=8
        )


def add_goals_layout_option_c(
    slide,
    main_lines: list[str],
    secondary_lines: list[str],
    *,
    top: float = 1.48,
) -> None:
    """Option C — Spin hand: three equal slot-ticket panels fanned + poker chip row."""
    usable_w = CONTENT_RIGHT - MARGIN_X
    cx = MARGIN_X + usable_w / 2

    cap = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(top), Inches(4.5), Inches(0.22))
    cap.text_frame.paragraphs[0].text = "MAIN GOALS · SPIN HAND"
    cap.text_frame.paragraphs[0].font.size = Pt(9)
    cap.text_frame.paragraphs[0].font.bold = True
    cap.text_frame.paragraphs[0].font.color.rgb = rgb(GOLD)

    card_w, card_h = 2.35, 2.75
    card_gap = 0.18
    card_y = top + 0.42
    row_w = 3 * card_w + 2 * card_gap
    row_left = cx - row_w / 2
    card_specs = (
        (-12, row_left, card_y, main_lines[0] if main_lines else "Main goal 1", 1, 0),
        (0, row_left + card_w + card_gap, card_y, main_lines[1] if len(main_lines) > 1 else "Main goal 2", 2, 1),
        (12, row_left + 2 * (card_w + card_gap), card_y, main_lines[2] if len(main_lines) > 2 else "Main goal 3", 3, 2),
    )
    for rot, x, y, line, num, idx in card_specs:
        outline = GOAL_MAIN_OUTLINES[idx % len(GOAL_MAIN_OUTLINES)]

        shadow = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.06),
            Inches(y + 0.06),
            Inches(card_w),
            Inches(card_h),
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = rgb(DARK_ALT)
        shadow.line.fill.background()
        shadow.rotation = float(rot)
        if shadow.adjustments:
            shadow.adjustments[0] = 0.06

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = rgb(DARK_CARD)
        _outline_shape(card, outline, pt=2.5)
        card.rotation = float(rot)
        if card.adjustments:
            card.adjustments[0] = 0.06

        stripe_h = 0.44
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(stripe_h)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = rgb(GOLD)
        stripe.line.fill.background()
        stripe.rotation = float(rot)

        ticket_lbl = slide.shapes.add_textbox(Inches(x), Inches(y + 0.08), Inches(card_w), Inches(0.28))
        ticket_lbl.text_frame.paragraphs[0].text = "SLOT TICKET"
        ticket_lbl.text_frame.paragraphs[0].font.size = Pt(8)
        ticket_lbl.text_frame.paragraphs[0].font.bold = True
        ticket_lbl.text_frame.paragraphs[0].font.color.rgb = rgb(BLACK)
        ticket_lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        ticket_lbl.rotation = float(rot)

        perf_x = x + card_w * 0.42
        perf = slide.shapes.add_shape(
            1,
            Inches(perf_x),
            Inches(y + stripe_h + 0.08),
            Inches(0.02),
            Inches(card_h - stripe_h - 0.16),
        )
        perf.fill.solid()
        perf.fill.fore_color.rgb = rgb(GOLD_MUTED)
        perf.line.fill.background()
        perf.rotation = float(rot)

        for notch_i in range(4):
            ny = y + stripe_h + 0.2 + notch_i * 0.52
            notch = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - 0.04),
                Inches(ny),
                Inches(0.1),
                Inches(0.1),
            )
            notch.fill.solid()
            notch.fill.fore_color.rgb = rgb(DARK)
            notch.line.color.rgb = rgb(outline)
            notch.line.width = Pt(1.0)
            notch.rotation = float(rot)

        nb = slide.shapes.add_textbox(Inches(x + card_w - 0.42), Inches(y + stripe_h + 0.12), Inches(0.32), Inches(0.3))
        nb.text_frame.paragraphs[0].text = str(num)
        nb.text_frame.paragraphs[0].font.size = Pt(16)
        nb.text_frame.paragraphs[0].font.bold = True
        nb.text_frame.paragraphs[0].font.color.rgb = rgb(outline)
        nb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        nb.rotation = float(rot)

        tb = slide.shapes.add_textbox(
            Inches(x + 0.18), Inches(y + stripe_h + 0.35), Inches(card_w - 0.36), Inches(card_h - stripe_h - 0.5)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(9)
        p.font.italic = _is_placeholder_line(line)
        p.font.color.rgb = rgb(MUTED if _is_placeholder_line(line) else BODY_TEXT)
        p.alignment = PP_ALIGN.CENTER
        tb.rotation = float(rot)

    sec_cap = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(4.95), Inches(4.5), Inches(0.2))
    sec_cap.text_frame.paragraphs[0].text = "SECONDARY GOALS · THE CHIPS"
    sec_cap.text_frame.paragraphs[0].font.size = Pt(8)
    sec_cap.text_frame.paragraphs[0].font.color.rgb = rgb(MUTED)

    chip_d = 1.05
    chip_gap = 0.1
    chip_y = 5.28
    for i in range(3):
        line = secondary_lines[i] if i < len(secondary_lines) else f"Secondary goal {i + 1}"
        outline = GOAL_SEC_OUTLINES[i % len(GOAL_SEC_OUTLINES)]
        px = _secondary_cluster_x(3, i, chip_d, chip_gap, cx)

        outer = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(px - 0.03), Inches(chip_y - 0.03), Inches(chip_d + 0.06), Inches(chip_d + 0.06)
        )
        outer.fill.solid()
        outer.fill.fore_color.rgb = rgb(outline)
        outer.line.fill.background()

        chip = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px), Inches(chip_y), Inches(chip_d), Inches(chip_d))
        chip.fill.solid()
        chip.fill.fore_color.rgb = rgb(DARK_VALUE)
        _outline_shape(chip, outline, pt=2.0)

        ring_d = chip_d * 0.72
        ring = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(px + (chip_d - ring_d) / 2),
            Inches(chip_y + (chip_d - ring_d) / 2),
            Inches(ring_d),
            Inches(ring_d),
        )
        ring.fill.background()
        ring.line.color.rgb = rgb(outline)
        ring.line.width = Pt(1.5)

        inner_d = chip_d * 0.38
        inner = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(px + (chip_d - inner_d) / 2),
            Inches(chip_y + (chip_d - inner_d) / 2),
            Inches(inner_d),
            Inches(inner_d),
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = rgb(DARK_CARD)
        inner.line.color.rgb = rgb(GOLD)
        inner.line.width = Pt(1.0)

        _goal_text_in_shape(slide, px + 0.04, chip_y + 0.22, chip_d - 0.08, 0.62, line, size=8)


def add_event_goals_split_view(
    slide,
    main_goals: list[tuple[str, str]],
    how_we_deliver: list[tuple[str, str]],
    *,
    top: float = 1.48,
) -> None:
    """Two-panel goals slide — Main Goals | Secondary Goals, slot-machine cabinet."""
    gap = 0.22
    arrow_w = 0.22
    usable_w = CONTENT_RIGHT - MARGIN_X
    panel_w = (usable_w - gap - arrow_w) / 2
    left_main = MARGIN_X
    left_secondary = MARGIN_X + panel_w + gap + arrow_w
    hdr_h = 0.42
    item_h = 1.02
    pad = 0.14
    panel_h = hdr_h + pad * 2 + 3 * item_h
    marquee_h = 0.26
    cabinet_pad = 0.1
    cabinet_top = top
    panels_top = cabinet_top + cabinet_pad + marquee_h + 0.08
    cabinet_h = panels_top - cabinet_top + panel_h + cabinet_pad

    cabinet = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN_X - 0.04),
        Inches(cabinet_top),
        Inches(usable_w + 0.08),
        Inches(cabinet_h),
    )
    cabinet.fill.solid()
    cabinet.fill.fore_color.rgb = rgb(DARK_ALT)
    _outline_shape(cabinet, GOLD, pt=2.5)
    if cabinet.adjustments:
        cabinet.adjustments[0] = 0.035

    marquee = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN_X + cabinet_pad),
        Inches(cabinet_top + cabinet_pad),
        Inches(usable_w - 2 * cabinet_pad),
        Inches(marquee_h),
    )
    marquee.fill.solid()
    marquee.fill.fore_color.rgb = rgb(GOLD)
    marquee.line.fill.background()
    if marquee.adjustments:
        marquee.adjustments[0] = 0.2

    marquee_txt = slide.shapes.add_textbox(
        Inches(MARGIN_X + cabinet_pad),
        Inches(cabinet_top + cabinet_pad + 0.04),
        Inches(usable_w - 2 * cabinet_pad),
        Inches(0.2),
    )
    marquee_txt.text_frame.paragraphs[0].text = "★  JACKPOTA · VEGAS 2026  ★"
    marquee_txt.text_frame.paragraphs[0].font.size = Pt(9)
    marquee_txt.text_frame.paragraphs[0].font.bold = True
    marquee_txt.text_frame.paragraphs[0].font.color.rgb = rgb(BLACK)
    marquee_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _draw_panel(
        left: float,
        title: str,
        items: list[tuple[str, str]],
        *,
        primary: bool,
        outlines: tuple[str, ...],
    ) -> None:
        y0 = panels_top
        glow = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left - 0.025),
            Inches(y0 - 0.025),
            Inches(panel_w + 0.05),
            Inches(panel_h + 0.05),
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = rgb(GOLD if primary else GOLD_MUTED)
        glow.line.fill.background()
        if glow.adjustments:
            glow.adjustments[0] = 0.04

        shell = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(y0),
            Inches(panel_w),
            Inches(panel_h),
        )
        shell.fill.solid()
        shell.fill.fore_color.rgb = rgb(DARK_CARD)
        _outline_shape(shell, GOLD if primary else GOLD_MUTED, pt=2.0 if primary else 1.5)
        if shell.adjustments:
            shell.adjustments[0] = 0.035

        hdr = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(y0),
            Inches(panel_w),
            Inches(hdr_h),
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = rgb(GOLD if primary else GOLD_MUTED)
        hdr.line.fill.background()
        if hdr.adjustments:
            hdr.adjustments[0] = 0.12

        hbox = slide.shapes.add_textbox(Inches(left), Inches(y0 + 0.08), Inches(panel_w), Inches(0.26))
        hp = hbox.text_frame.paragraphs[0]
        hp.text = title
        hp.font.size = Pt(10)
        hp.font.bold = True
        hp.font.color.rgb = rgb(BLACK if primary else BODY_TEXT)
        hp.alignment = PP_ALIGN.CENTER

        body_top = y0 + hdr_h + pad
        for i, (line, _) in enumerate(items[:3]):
            outline = outlines[i % len(outlines)]
            y = body_top + i * item_h
            row = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + 0.1),
                Inches(y + 0.06),
                Inches(panel_w - 0.2),
                Inches(item_h - 0.12),
            )
            row.fill.solid()
            row.fill.fore_color.rgb = rgb(DARK_VALUE)
            _outline_shape(row, outline, pt=1.5)
            if row.adjustments:
                row.adjustments[0] = 0.12

            cx = left + 0.28
            cy = y + item_h * 0.42
            disc_d = Inches(0.34)
            disc = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx) - disc_d / 2,
                Inches(cy) - disc_d / 2,
                disc_d,
                disc_d,
            )
            disc.fill.solid()
            disc.fill.fore_color.rgb = rgb(DARK_CARD)
            disc.line.color.rgb = rgb(outline)
            disc.line.width = Pt(2.0)

            nbox = slide.shapes.add_textbox(
                Inches(cx) - disc_d / 2,
                Inches(cy) - disc_d / 2 + Inches(0.04),
                disc_d,
                disc_d,
            )
            np = nbox.text_frame.paragraphs[0]
            np.text = str(i + 1)
            np.font.size = Pt(12)
            np.font.bold = True
            np.font.color.rgb = rgb(GOLD)
            np.alignment = PP_ALIGN.CENTER

            if primary:
                sym_x = left + 0.46
                sym_w = 0.06
                for j in range(3):
                    sym = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(sym_x + j * 0.08),
                        Inches(y + item_h * 0.28),
                        Inches(sym_w),
                        Inches(0.22),
                    )
                    sym.fill.solid()
                    sym.fill.fore_color.rgb = rgb(outline if j == 1 else GOLD_MUTED)
                    sym.line.color.rgb = rgb(GOLD_MUTED)
                    sym.line.width = Pt(0.5)

            if i < min(len(items), 3) - 1:
                stripe = slide.shapes.add_shape(
                    1,
                    Inches(cx - 0.01),
                    Inches(cy + float(disc_d / Inches(1)) / 2 - 0.02),
                    Inches(0.02),
                    Inches(item_h - float(disc_d / Inches(1)) + 0.06),
                )
                stripe.fill.solid()
                stripe.fill.fore_color.rgb = rgb(outline)
                stripe.line.fill.background()

            text_x = left + (0.72 if primary else 0.52)
            text_w = panel_w - (0.82 if primary else 0.62)
            tbox = slide.shapes.add_textbox(Inches(text_x), Inches(y), Inches(text_w), Inches(item_h))
            ttf = tbox.text_frame
            ttf.word_wrap = True
            ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
            if " — " in line:
                head, tail = line.split(" — ", 1)
                tp = ttf.paragraphs[0]
                tp.text = head
                tp.font.size = Pt(10)
                tp.font.bold = True
                tp.font.color.rgb = rgb(GOLD)
                if tail:
                    sp = ttf.add_paragraph()
                    sp.text = tail
                    sp.font.size = Pt(8)
                    sp.font.color.rgb = rgb(BODY_TEXT)
            else:
                tp = ttf.paragraphs[0]
                tp.text = line
                tp.font.size = Pt(10 if len(line) < 42 else 9)
                tp.font.bold = True
                tp.font.color.rgb = rgb(BODY_TEXT)

    _draw_panel(left_main, "MAIN GOALS · THE REELS", main_goals, primary=True, outlines=GOAL_MAIN_OUTLINES)
    _draw_panel(
        left_secondary,
        "SECONDARY GOALS · BONUS LINE",
        how_we_deliver,
        primary=False,
        outlines=GOAL_SEC_OUTLINES,
    )

    mid_x = MARGIN_X + usable_w / 2
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(mid_x - arrow_w / 2),
        Inches(panels_top + panel_h / 2 - 0.1),
        Inches(arrow_w),
        Inches(0.2),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(GOLD)
    arrow.line.fill.background()


def _normalize_kpi_label(label: str) -> str:
    text = re.sub(r"\s+", " ", label.replace("\n", " ").strip())
    text = text.replace("DAYS FROM FTP", "SENIORITY GROUP")
    return text


def _split_kpi_label(label: str) -> tuple[str, str]:
    """Split metric label into standout headline + subline (preserves Excel newlines)."""
    raw = label.replace("\r", "").strip()
    parts = [re.sub(r"\s+", " ", p.strip()) for p in raw.split("\n") if p.strip()]
    if len(parts) >= 2:
        head = parts[0].replace("DAYS FROM FTP", "SENIORITY GROUP")
        sub = " ".join(parts[1:]).replace("DAYS FROM FTP", "SENIORITY GROUP")
        return head.upper(), sub.upper()
    text = _normalize_kpi_label(label)
    words = text.split()
    if text.upper().startswith("TOP STATE"):
        return "TOP STATE", " ".join(words[2:]).upper() if len(words) > 2 else "BY COUNT"
    if text.upper().startswith("AVG "):
        return " ".join(words[:3]).upper(), " ".join(words[3:]).upper()
    if text.upper().startswith("MEDIAN "):
        return "MEDIAN", " ".join(words[1:]).upper()
    if len(words) <= 2:
        return text.upper(), ""
    mid = max(1, len(words) // 2)
    return " ".join(words[:mid]).upper(), " ".join(words[mid:]).upper()


def _fit_circle_value(val: str, circle_d_inches: float, n_cards: int) -> tuple[str, int]:
    """Format value for circle interior — auto size and optional two-line split."""
    display = val

    def _chars_per_line(pt: int) -> int:
        d_pt = circle_d_inches * 72
        return max(3, int(d_pt * 0.52 / (pt * 0.48)))

    def _fits(text: str, pt: int) -> bool:
        lines = text.split("\n")
        cap = _chars_per_line(pt)
        return len(lines) <= 2 and all(len(line) <= cap for line in lines)

    for pt in range(12, 6, -1):
        use_pt = max(7, pt - (1 if n_cards > 6 else 0))
        if _fits(display, use_pt):
            return display, use_pt
        if "," in display and "\n" not in display:
            parts = display.rsplit(",", 1)
            two = f"{parts[0]},\n{parts[1]}"
            if _fits(two, use_pt):
                return two, use_pt
        if "\n" not in display:
            mid = (len(display) + 1) // 2
            two = f"{display[:mid]}\n{display[mid:]}"
            if _fits(two, use_pt):
                return two, use_pt
    return display, 7


def _add_circle_medallion(slide, cx, cy, diameter, *, ring_pt: float = 2.0):
    """Gold-ring medallion — stays within diameter bounds (no outer glow)."""
    r = diameter / 2
    x, y = cx - r, cy - r
    disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, diameter, diameter)
    disc.fill.solid()
    disc.fill.fore_color.rgb = rgb(DARK_VALUE)
    disc.line.color.rgb = rgb(GOLD)
    disc.line.width = Pt(ring_pt)
    return disc


def cohort_two_chart_layout(content_top: float) -> tuple[float, float, float, float, float]:
    """Return left1, left2, chart_top, chart_w, chart_h (inches) within slide bounds."""
    gap = 0.26
    section_h = 0.36
    chart_top = content_top + 0.1
    usable_w = CONTENT_RIGHT - MARGIN_X
    chart_w = (usable_w - gap) / 2
    left1 = MARGIN_X
    left2 = MARGIN_X + chart_w + gap
    max_h = CONTENT_BOTTOM - chart_top - section_h
    chart_h = min(max(max_h, 2.6), 3.25)
    return left1, left2, chart_top, chart_w, chart_h


def add_kpi_row(slide, top: float, cards: list[tuple[str, str]]) -> float:
    """KPI strip — gold medallion circles; sized to fit slide width."""
    n = len(cards)
    if n == 0:
        return 0.0

    margin = Inches(MARGIN_X)
    total_w = Inches(CONTENT_RIGHT - MARGIN_X)
    gap = Inches(0.05)
    slot_w = (total_w - gap * (n - 1)) / n
    y_top = Inches(top)

    # Scale circle to slot — cap size so 7 metrics fit without horizontal bleed
    max_d_in = 0.66 if n >= 7 else (0.72 if n >= 5 else 0.78)
    circle_d = min(slot_w * 0.74, Inches(max_d_in))
    circle_cy = y_top + circle_d / 2 + Inches(0.02)
    label_top = circle_cy + circle_d / 2 + Inches(0.06)
    label_h = Inches(0.36)
    head_pt = 6 if n > 6 else 7
    sub_pt = 5 if n > 6 else 6
    circle_d_in = float(circle_d / Inches(1))

    for i, (val, lbl) in enumerate(cards):
        slot_x = margin + i * (slot_w + gap)
        cx = slot_x + slot_w / 2

        _add_circle_medallion(slide, cx, circle_cy, circle_d)

        display_val = val if val not in ("0", "0%", "0.0%") else "—"
        extra_sub = ""
        if len(display_val) > 12 and "·" in display_val:
            pct_m = re.search(r"\((\d+%)\)", display_val)
            if pct_m:
                extra_sub = display_val.split("·")[0].strip()
                display_val = pct_m.group(1)

        display_val, val_pt = _fit_circle_value(display_val, circle_d_in, n)
        inset = circle_d * 0.16
        vbox = slide.shapes.add_textbox(
            cx - circle_d / 2 + inset,
            circle_cy - circle_d / 2 + inset,
            circle_d - inset * 2,
            circle_d - inset * 2,
        )
        vtf = vbox.text_frame
        vtf.word_wrap = True
        vtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        vtf.margin_left = 0
        vtf.margin_right = 0
        vtf.margin_top = 0
        vtf.margin_bottom = 0
        for line_i, line in enumerate(display_val.split("\n")):
            vp = vtf.paragraphs[0] if line_i == 0 else vtf.add_paragraph()
            vp.text = line
            vp.font.size = Pt(val_pt)
            vp.font.bold = True
            vp.font.color.rgb = rgb(GOLD)
            vp.alignment = PP_ALIGN.CENTER
            vp.space_before = Pt(0)
            vp.space_after = Pt(0)

        headline, subline = _split_kpi_label(lbl)
        if extra_sub and not subline:
            subline = extra_sub
            extra_sub = ""

        lbox = slide.shapes.add_textbox(
            slot_x,
            label_top,
            slot_w,
            label_h,
        )
        ltf = lbox.text_frame
        ltf.word_wrap = True
        ltf.vertical_anchor = MSO_ANCHOR.TOP
        ltf.margin_left = Pt(1)
        ltf.margin_right = Pt(1)
        hp = ltf.paragraphs[0]
        hp.text = headline
        hp.font.size = Pt(head_pt)
        hp.font.bold = True
        hp.font.color.rgb = rgb(GOLD)
        hp.alignment = PP_ALIGN.CENTER
        hp.space_after = Pt(1)

        if subline:
            sp = ltf.add_paragraph()
            sp.text = subline
            sp.font.size = Pt(sub_pt)
            sp.font.bold = True
            sp.font.color.rgb = rgb(BODY_TEXT)
            sp.alignment = PP_ALIGN.CENTER
        elif extra_sub:
            sp = ltf.add_paragraph()
            sp.text = extra_sub
            sp.font.size = Pt(sub_pt)
            sp.font.color.rgb = rgb(MUTED)
            sp.alignment = PP_ALIGN.CENTER

    bottom = float(label_top / Inches(1)) + float(label_h / Inches(1))
    return bottom - top + 0.04


def _chart_series_colors(chart, n: int) -> None:
    for i in range(n):
        try:
            pt = chart.series[0].points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = rgb(CHART_COLORS[i % len(CHART_COLORS)])
        except (IndexError, AttributeError):
            break


def _ensure_child(parent, tag: str, val: str | None = None):
    ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    el = parent.find(qn(f"c:{tag}"))
    if el is None:
        from pptx.oxml import parse_xml

        el = parse_xml(f'<c:{tag} xmlns:c="{ns}"/>')
        parent.append(el)
    if val is not None:
        el.set("val", val)
    return el


def _configure_pie_data_labels(chart, n_slices: int) -> None:
    """Category name + percent inside each slice; hide raw counts to avoid stray zeros."""
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(6 if n_slices > 6 else 7)
    dl.font.color.rgb = rgb(BODY_TEXT)
    dl.number_format = "0%"

    d_lbls = plot._element.find(qn("c:dLbls"))
    if d_lbls is None:
        return
    for tag, val in (
        ("showLegendKey", "0"),
        ("showCatName", "1"),
        ("showVal", "0"),
        ("showPercent", "1"),
        ("showSerName", "0"),
        ("showLeaderLines", "0"),
        ("dLblPos", "bestFit"),
    ):
        _ensure_child(d_lbls, tag, val)
    sep = d_lbls.find(qn("c:separator"))
    if sep is None:
        from pptx.oxml import parse_xml

        sep = parse_xml(
            '<c:separator xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"/>'
        )
        d_lbls.append(sep)
    sep.text = "\n"


def add_section_label(slide, left: float, top: float, width: float, title: str) -> None:
    bar_h = 0.3
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(bar_h),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(GOLD_MUTED)
    bar.line.fill.background()
    if bar.adjustments:
        bar.adjustments[0] = 0.15

    box = slide.shapes.add_textbox(
        Inches(left + 0.14), Inches(top + 0.05), Inches(width - 0.2), Inches(bar_h - 0.08)
    )
    p = box.text_frame.paragraphs[0]
    p.text = title.strip()
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = rgb(GOLD)


def _filter_chart_series(
    categories: list[str], values: list[float]
) -> tuple[list[str], list[float]]:
    """Drop zero-count slices so labels and legends stay clean."""
    cats: list[str] = []
    vals: list[float] = []
    for cat, val in zip(categories, values):
        try:
            v = float(val)
        except (TypeError, ValueError):
            continue
        if v <= 0:
            continue
        cats.append(str(cat))
        vals.append(v)
    return cats, vals


def add_pie_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    categories: list[str],
    values: list[float],
) -> None:
    categories, values = _filter_chart_series(categories, values)
    if not categories or not values:
        return
    add_section_label(slide, left, top, width, title)
    chart_top = top + 0.34
    chart_h = max(height - 0.34, 2.4)
    data = CategoryChartData()
    data.categories = categories
    data.add_series("", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(left),
        Inches(chart_top),
        Inches(width),
        Inches(chart_h),
        data,
    ).chart
    chart.has_legend = False
    _configure_pie_data_labels(chart, len(categories))
    _chart_series_colors(chart, len(categories))


def add_bar_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    categories: list[str],
    values: list[float],
    *,
    horizontal: bool = True,
) -> None:
    if not categories or not values:
        return
    add_section_label(slide, left, top, width, title)
    chart_top = top + 0.32
    chart_h = height - 0.32
    data = CategoryChartData()
    data.categories = categories
    data.add_series("", values)
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart = slide.shapes.add_chart(
        chart_type,
        Inches(left),
        Inches(chart_top),
        Inches(width),
        Inches(chart_h),
        data,
    ).chart
    chart.has_legend = False
    if chart.series:
        for i, _ in enumerate(categories):
            try:
                pt = chart.series[0].points[i]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = rgb(CHART_COLORS[i % len(CHART_COLORS)])
            except (IndexError, AttributeError):
                break
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.color.rgb = rgb(BODY_TEXT)


def add_panel_table(
    slide,
    left: float,
    top: float,
    width: float,
    title: str,
    headers: tuple[str, ...],
    rows: list[tuple],
    *,
    pct_col: int | None = None,
) -> None:
    add_section_label(slide, left, top, width, title)
    if not rows:
        return
    tbl_top = top + 0.34
    row_h = 0.28
    nrows = 1 + len(rows)
    table = slide.shapes.add_table(
        nrows, len(headers), Inches(left), Inches(tbl_top), Inches(width), Inches(row_h * nrows)
    ).table

    for j, h in enumerate(headers):
        style_table_header(table.cell(0, j), h)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            text = val
            if pct_col is not None and j == pct_col and isinstance(val, (int, float)):
                text = f"{val:.0%}" if val <= 1 else str(val)
            style_table_body(table.cell(i + 1, j), str(text) if text is not None else "", alt_row=i % 2 == 1)


def _schedule_item_style(text: str) -> tuple[str, bool, bool]:
    """Return display text, highlight (gold), optional (muted italic)."""
    display = text
    if "Gala" in text and "13:00" in text:
        display = "Gala · 13:00–16:00"
    highlight = "gala" in display.lower()
    optional = "optional" in display.lower()
    if "no pricing" in display.lower():
        display = display.replace("(no pricing)", "").replace("(No pricing)", "").strip()
    return display, highlight, optional


def add_day_cards_schedule(
    slide,
    schedule_days: tuple[tuple[str, ...], ...],
    *,
    footer: str = "",
    top: float | None = None,
) -> None:
    """Three timeline day cards — rounded panels, gold headers, clean vertical flow."""
    margin_x = MARGIN_X
    gap = 0.14
    n_days = len(schedule_days)
    usable_w = CONTENT_RIGHT - margin_x
    block_w = usable_w * 0.78
    block_left = margin_x + (usable_w - block_w) / 2
    card_w = (block_w - gap * (n_days - 1)) / n_days

    max_items = max(len(d) for d in schedule_days)
    hdr_h = 0.4
    item_h = 0.33
    body_pad_t = 0.16
    body_pad_b = 0.18
    body_pad_x = 0.16
    card_body_h = body_pad_t + body_pad_b + max_items * item_h
    card_h = hdr_h + card_body_h

    footer_h = 0.32 if footer else 0.0
    area_top = 1.52
    area_bottom = CONTENT_BOTTOM - footer_h - 0.12
    if top is None:
        top = area_top + max(0.0, (area_bottom - area_top - card_h) * 0.35)

    day_labels = ("DAY 1", "DAY 2", "DAY 3")
    day_tags = ("Arrival", "Main Event", "Departure")

    for j, day_items in enumerate(schedule_days):
        left = block_left + j * (card_w + gap)

        # Card shell — rounded dark panel
        shell = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(card_w),
            Inches(card_h),
        )
        shell.fill.solid()
        shell.fill.fore_color.rgb = rgb(DARK_CARD)
        shell.line.color.rgb = rgb(GOLD_MUTED)
        shell.line.width = Pt(0.75)
        if shell.adjustments:
            shell.adjustments[0] = 0.04

        # Gold header band
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(card_w),
            Inches(hdr_h),
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = rgb(GOLD)
        hdr.line.fill.background()
        if hdr.adjustments:
            hdr.adjustments[0] = 0.12

        day_lbl = day_labels[j] if j < len(day_labels) else f"DAY {j + 1}"
        day_tag = day_tags[j] if j < len(day_tags) else ""
        hbox = slide.shapes.add_textbox(
            Inches(left), Inches(top + 0.06), Inches(card_w), Inches(0.22)
        )
        hp = hbox.text_frame.paragraphs[0]
        hp.text = day_lbl
        hp.font.size = Pt(13)
        hp.font.bold = True
        hp.font.color.rgb = rgb(BLACK)
        hp.alignment = PP_ALIGN.CENTER

        if day_tag:
            tbox = slide.shapes.add_textbox(
                Inches(left), Inches(top + 0.28), Inches(card_w), Inches(0.18)
            )
            tp = tbox.text_frame.paragraphs[0]
            tp.text = day_tag
            tp.font.size = Pt(7)
            tp.font.bold = True
            tp.font.color.rgb = rgb(BLACK)
            tp.alignment = PP_ALIGN.CENTER

        body_top = top + hdr_h
        # Left gold accent stripe
        stripe = slide.shapes.add_shape(
            1,
            Inches(left + 0.06),
            Inches(body_top + 0.12),
            Inches(0.04),
            Inches(card_body_h - 0.24),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = rgb(GOLD)
        stripe.line.fill.background()

        dot_x = left + 0.22
        text_x = left + body_pad_x + 0.28
        text_w = card_w - body_pad_x - 0.38
        row_top = body_top + body_pad_t

        for i, raw in enumerate(day_items):
            text, highlight, optional = _schedule_item_style(raw)
            y = row_top + i * item_h
            dot_y = y + item_h * 0.32

            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(dot_x),
                Inches(dot_y),
                Inches(0.09),
                Inches(0.09),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(GOLD if highlight else GOLD_MUTED)
            dot.line.fill.background()

            if i < len(day_items) - 1:
                line = slide.shapes.add_shape(
                    1,
                    Inches(dot_x + 0.04),
                    Inches(dot_y + 0.09),
                    Inches(0.01),
                    Inches(item_h - 0.05),
                )
                line.fill.solid()
                line.fill.fore_color.rgb = rgb(GOLD_MUTED)
                line.line.fill.background()

            ibox = slide.shapes.add_textbox(
                Inches(text_x), Inches(y), Inches(text_w), Inches(item_h)
            )
            itf = ibox.text_frame
            itf.word_wrap = True
            itf.vertical_anchor = MSO_ANCHOR.MIDDLE
            itf.margin_left = 0
            itf.margin_right = 0
            ip = itf.paragraphs[0]
            ip.text = text
            ip.font.size = Pt(10 if highlight else 9)
            ip.font.bold = highlight
            ip.font.color.rgb = rgb(GOLD if highlight else BODY_TEXT)
            if optional:
                ip.font.italic = True
                ip.font.color.rgb = rgb(MUTED)
                ip.font.bold = False

    if footer:
        fbox = slide.shapes.add_textbox(
            Inches(margin_x),
            Inches(CONTENT_BOTTOM - footer_h),
            Inches(usable_w),
            Inches(footer_h),
        )
        fp = fbox.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(11)
        fp.font.italic = True
        fp.font.color.rgb = rgb(GOLD)
        fp.alignment = PP_ALIGN.CENTER


def add_schedule_table(slide, schedule_days: tuple[tuple[str, ...], ...], *, footer: str = "") -> None:
    """Legacy table layout — delegates to day cards."""
    add_day_cards_schedule(slide, schedule_days, footer=footer)


def delete_slide(prs: Presentation, index: int) -> None:
    r_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def delete_first_slides(prs: Presentation, count: int) -> int:
    removed = 0
    for _ in range(min(count, len(prs.slides))):
        delete_slide(prs, 0)
        removed += 1
    return removed


def prepend_slides(prs: Presentation, count: int) -> None:
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_slides = slides[-count:]
    for s in new_slides:
        xml_slides.remove(s)
    for i, s in enumerate(new_slides):
        xml_slides.insert(i, s)


def replace_first_slides(prs: Presentation, count: int, build_fn) -> None:
    """Delete first `count` slides, append new ones via build_fn, move them to front."""
    delete_first_slides(prs, count)
    before = len(prs.slides)
    build_fn()
    added = len(prs.slides) - before
    if added:
        prepend_slides(prs, added)


def remove_duplicate_schedule_slides(prs: Presentation, keep_through_index: int = 2) -> int:
    to_remove: list[int] = []
    for idx, slide in enumerate(prs.slides):
        if idx <= keep_through_index:
            continue
        texts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    texts.append(t)
        combined = " ".join(texts)
        if "Schedule" in combined and "Top 50 Elite" in combined:
            to_remove.append(idx)

    removed = 0
    for idx in reversed(to_remove):
        delete_slide(prs, idx)
        removed += 1
    return removed
