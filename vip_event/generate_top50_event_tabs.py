"""
VIP Event Top 50 — Schedule & Focus Points tabs (Excel + optional PowerPoint).

Adds/replaces only Schedule and Focus Points sheets in the Top 50 workbook.
Does not modify the player roster tab or any Top 30 artifacts.

Usage:
  python vip_event/generate_top50_event_tabs.py
  python vip_event/generate_top50_event_tabs.py --source "path/to/VIP Event - Top 50 Players .xlsx"
  python vip_event/generate_top50_event_tabs.py --pptx "path/to/VIP Event - Vegas 2026.pptx"
"""
from __future__ import annotations

import argparse
import shutil
import sys
from collections import Counter
from dataclasses import dataclass
from pathlib import Path

from openpyxl import load_workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx.util import Inches, Pt

MODULE_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(MODULE_DIR))
from pptx_theme import (  # noqa: E402
    TITLE_MAIN,
    add_day_cards_schedule,
    add_event_goals_split_view,
    add_executive_bullets,
    add_title_band,
    load_template,
    set_speaker_notes,
)
from event_playbook_content import (  # noqa: E402
    EVENT_GOALS,
    EVENT_GOALS_DELIVERY,
    EVENT_GOALS_MAIN,
    EVENT_GOALS_SUBTITLE,
    FOCUS_POINTS,
    FOCUS_POINTS_SUBTITLE,
)

EXPORT_DIR = MODULE_DIR / "exports"

DEFAULT_SOURCE = Path(
    r"c:\Users\Owner\OneDrive - Silver Social Games\Desktop\VIP\VIP Event"
    r"\VIP Event - Top 50 Players .xlsx"
)
FALLBACK_SOURCE = MODULE_DIR / "data" / "vip-event-top50-players.xlsx"
DEFAULT_PPTX = Path(
    r"c:\Users\Owner\OneDrive - Silver Social Games\Desktop\VIP\VIP Event"
    r"\VIP Event - Vegas 2026.pptx"
)

ROSTER_SHEET = "VIP Event - Top 50 Players "
SCHEDULE_SHEET = "Schedule"
FOCUS_SHEET = "Focus Points"

DATA_START = 9
DATA_END = 58
COL_STATE, COL_AGENT = 5, 6
COL_NP_LT, COL_NP_30 = 7, 8

# 3-day itinerary from event plan
SCHEDULE_DAYS: tuple[tuple[str, ...], ...] = (
    (
        "Pick-up from airport",
        "Gathering",
        "Check-in",
        "Free time",
        "Dinner",
        "Sphere (optional)",
    ),
    (
        "Self breakfast (no pricing)",
        "Gala (3 hours: 13:00-16:00)",
        "Free time",
        "Dinner",
        "Show (Music)",
    ),
    (
        "Brunch + closure",
        "Pick up from hotel",
    ),
)

# Executive palette (matches generate_top50_management_brief.py)
NAVY = "FF1B2A4A"
NAVY_MID = "FF2C3E6E"
ORANGE = "FFED7D31"
ORANGE_LIGHT = "FFFFF4E8"
HEADER_BG = "FFE8EEF5"
INK = "FF1A1A2E"
WHITE = "FFFFFFFF"
SURFACE_ALT = "FFF0F4F8"
MUTED = "FF6B7C93"
RULE = "FFE2E8F0"

SIDE_THIN = Side(style="thin", color=RULE)
BORDER = Border(left=SIDE_THIN, right=SIDE_THIN, top=SIDE_THIN, bottom=SIDE_THIN)

FONT_TITLE = Font(name="Calibri", size=28, bold=True, color=WHITE)
FONT_SUB = Font(name="Calibri", size=14, color="FFB8C9E0")
FONT_DAY_HDR = Font(name="Calibri", size=13, bold=True, color=WHITE)
FONT_TH = Font(name="Calibri", size=12, bold=True, color=INK[2:])
FONT_BODY = Font(name="Calibri", size=12, color=INK[2:])
FONT_FOCUS_NUM = Font(name="Calibri", size=14, bold=True, color=ORANGE[2:])
FONT_FOCUS_HEAD = Font(name="Calibri", size=12, bold=True, color=INK[2:])
FONT_FOCUS_BODY = Font(name="Calibri", size=11, color=MUTED[2:])
FONT_FOOT = Font(name="Calibri", size=10, italic=True, color=MUTED[2:])

LAST_COL = 3


@dataclass
class Top50Stats:
    player_count: int
    sum_np_lt: float
    sum_np_30d: float
    top_state: str
    top_state_n: int
    top_state_pct: float
    agent_counts: dict[str, int]
    low_momentum_n: int  # LT >= 50k and 30d < 5k


def _resolve_source(path: Path | None) -> Path:
    if path and path.exists():
        return path
    if DEFAULT_SOURCE.exists():
        return DEFAULT_SOURCE
    if FALLBACK_SOURCE.exists():
        return FALLBACK_SOURCE
    raise FileNotFoundError(
        f"Top 50 workbook not found. Tried: {DEFAULT_SOURCE}, {FALLBACK_SOURCE}"
    )


def _open_workbook(preferred: Path) -> tuple:
    """Load workbook; fall back to data copy if OneDrive file is locked open."""
    candidates = [preferred]
    if preferred == DEFAULT_SOURCE and FALLBACK_SOURCE not in candidates:
        candidates.append(FALLBACK_SOURCE)
    elif preferred != FALLBACK_SOURCE and FALLBACK_SOURCE.exists():
        candidates.append(FALLBACK_SOURCE)

    last_err: Exception | None = None
    for path in candidates:
        try:
            return load_workbook(path), path
        except PermissionError as e:
            last_err = e
    raise PermissionError(
        f"Cannot open workbook (close it in Excel): {preferred}"
    ) from last_err


def _num(val) -> float:
    if val is None:
        return 0.0
    try:
        return float(val)
    except (TypeError, ValueError):
        return 0.0


def read_top50_stats(ws) -> Top50Stats:
    players: list[dict] = []
    for row in range(DATA_START, DATA_END + 1):
        aid = ws.cell(row, 2).value
        if aid is None:
            continue
        players.append(
            {
                "state": str(ws.cell(row, COL_STATE).value or "—"),
                "agent": str(ws.cell(row, COL_AGENT).value or "—"),
                "np_lt": _num(ws.cell(row, COL_NP_LT).value),
                "np_30d": _num(ws.cell(row, COL_NP_30).value),
            }
        )

    state_ctr = Counter(p["state"] for p in players)
    agent_ctr = Counter(p["agent"] for p in players)
    top_state, top_n = state_ctr.most_common(1)[0] if state_ctr else ("—", 0)
    n = len(players) or 1

    return Top50Stats(
        player_count=len(players),
        sum_np_lt=sum(p["np_lt"] for p in players),
        sum_np_30d=sum(p["np_30d"] for p in players),
        top_state=top_state,
        top_state_n=top_n,
        top_state_pct=top_n / n,
        agent_counts=dict(agent_ctr.most_common()),
        low_momentum_n=sum(1 for p in players if p["np_lt"] >= 50_000 and p["np_30d"] < 5_000),
    )


def set_cell(ws, row, col, value=None, font=None, fill=None, align=None, border=None):
    c = ws.cell(row, col, value)
    if font:
        c.font = font
    if fill:
        c.fill = fill
    if align:
        c.alignment = align
    if border:
        c.border = border
    return c


def merge_style(ws, r1, c1, r2, c2, value=None, font=None, fill=None, align=None):
    ws.merge_cells(start_row=r1, start_column=c1, end_row=r2, end_column=c2)
    set_cell(ws, r1, c1, value, font, fill, align)


def _replace_sheet(wb, name: str):
    if name in wb.sheetnames:
        del wb[name]
    return wb.create_sheet(name)


def build_focus_points(_stats: Top50Stats | None = None) -> list[tuple[str, str]]:
    """Operational focus points — (slide_line, speaker_note)."""
    return list(FOCUS_POINTS)


def build_schedule_sheet(wb) -> None:
    ws = _replace_sheet(wb, SCHEDULE_SHEET)
    ws.sheet_view.showGridLines = False

    navy = PatternFill("solid", fgColor=NAVY)
    navy_mid = PatternFill("solid", fgColor=NAVY_MID)
    orange = PatternFill("solid", fgColor=ORANGE)
    hdr_fill = PatternFill("solid", fgColor=HEADER_BG)
    alt_fill = PatternFill("solid", fgColor=SURFACE_ALT)
    white = PatternFill("solid", fgColor=WHITE)

    merge_style(
        ws, 1, 1, 2, LAST_COL,
        "VIP Event · Vegas 2026 · Schedule",
        FONT_TITLE, navy, Alignment(horizontal="left", vertical="center", indent=1),
    )
    merge_style(
        ws, 3, 1, 3, LAST_COL,
        "Top 50 Elite invitees · 3-day program",
        FONT_SUB, navy_mid, Alignment(horizontal="left", vertical="center", indent=1),
    )
    ws.row_dimensions[1].height = 38
    ws.row_dimensions[2].height = 6
    ws.row_dimensions[3].height = 24
    for c in range(1, LAST_COL + 1):
        set_cell(ws, 4, c, fill=orange)

    day_align = Alignment(horizontal="center", vertical="center")
    body_align = Alignment(horizontal="left", vertical="top", wrap_text=True, indent=1)

    for col, label in enumerate(("Day 1", "Day 2", "Day 3"), start=1):
        set_cell(ws, 5, col, label, FONT_DAY_HDR, orange, day_align, BORDER)
    ws.row_dimensions[5].height = 28

    max_rows = max(len(d) for d in SCHEDULE_DAYS)
    for i in range(max_rows):
        row = 6 + i
        row_fill = alt_fill if i % 2 else white
        ws.row_dimensions[row].height = 26
        for col, day_items in enumerate(SCHEDULE_DAYS, start=1):
            text = day_items[i] if i < len(day_items) else None
            set_cell(ws, row, col, text, FONT_BODY, row_fill, body_align, BORDER)

    foot = 6 + max_rows + 2
    merge_style(
        ws, foot, 1, foot, LAST_COL,
        "Notes: Day 2 breakfast is self-serve (not priced). Sphere on Day 1 is optional.",
        FONT_FOOT, white, Alignment(horizontal="left", indent=1),
    )

    for col, w in enumerate((36, 36, 36), start=1):
        ws.column_dimensions[get_column_letter(col)].width = w

    ws.page_setup.orientation = "landscape"
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 1
    ws.print_area = f"A1:{get_column_letter(LAST_COL)}{foot}"


def build_focus_points_sheet(wb, stats: Top50Stats) -> list[tuple[str, str]]:
    points = build_focus_points(stats)
    ws = _replace_sheet(wb, FOCUS_SHEET)
    ws.sheet_view.showGridLines = False

    navy = PatternFill("solid", fgColor=NAVY)
    navy_mid = PatternFill("solid", fgColor=NAVY_MID)
    orange = PatternFill("solid", fgColor=ORANGE)
    orange_light = PatternFill("solid", fgColor=ORANGE_LIGHT)
    white = PatternFill("solid", fgColor=WHITE)
    alt_fill = PatternFill("solid", fgColor=SURFACE_ALT)

    merge_style(
        ws, 1, 1, 2, 4,
        "VIP Event · Vegas 2026 · Focus Points",
        FONT_TITLE, navy, Alignment(horizontal="left", vertical="center", indent=1),
    )
    merge_style(
        ws, 3, 1, 3, 4,
        "Top 50 Elite · operational prep (slide line + speaker notes)",
        FONT_SUB, navy_mid, Alignment(horizontal="left", vertical="center", indent=1),
    )
    ws.row_dimensions[1].height = 38
    ws.row_dimensions[2].height = 6
    ws.row_dimensions[3].height = 24
    for c in range(1, 5):
        set_cell(ws, 4, c, fill=orange)

    hdr_row = 5
    for col, label in ((1, "#"), (2, "Slide"), (3, "Notes")):
        if col == 3:
            merge_style(ws, hdr_row, col, hdr_row, 4, label, FONT_TH, PatternFill("solid", fgColor=HEADER_BG),
                        Alignment(horizontal="left", vertical="center", indent=1))
        else:
            set_cell(ws, hdr_row, col, label, FONT_TH, PatternFill("solid", fgColor=HEADER_BG),
                     Alignment(horizontal="center" if col == 1 else "left", vertical="center", indent=1), BORDER)
    ws.row_dimensions[hdr_row].height = 22

    for i, (slide_line, note) in enumerate(points):
        row = hdr_row + 1 + i
        row_fill = alt_fill if i % 2 else white
        set_cell(ws, row, 1, i + 1, FONT_FOCUS_NUM, orange_light,
                 Alignment(horizontal="center", vertical="top"), BORDER)
        set_cell(ws, row, 2, slide_line, FONT_FOCUS_HEAD, row_fill,
                 Alignment(horizontal="left", vertical="top", wrap_text=True, indent=1), BORDER)
        merge_style(ws, row, 3, row, 4, note, FONT_FOCUS_BODY, row_fill,
                    Alignment(horizontal="left", vertical="top", wrap_text=True, indent=1))
        ws.row_dimensions[row].height = 44

    foot = hdr_row + len(points) + 2
    merge_style(
        ws, foot, 1, foot, 4,
        "Source: VIP Event - Top 50 Players roster tab · Managed Elite book",
        FONT_FOOT, white, Alignment(horizontal="left", indent=1),
    )

    ws.column_dimensions["A"].width = 5
    ws.column_dimensions["B"].width = 22
    ws.column_dimensions["C"].width = 55
    ws.column_dimensions["D"].width = 12

    ws.page_setup.orientation = "landscape"
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 1
    ws.print_area = f"A1:D{foot}"
    return points


def write_pptx(
    schedule_days: tuple[tuple[str, ...], ...],
    focus_points: list[tuple[str, str]],
    template: Path | None,
    output: Path,
) -> Path:
    prs = load_template(str(template) if template and template.exists() else None)

    add_title_band(prs.slides.add_slide(prs.slide_layouts[6]), TITLE_MAIN, "Schedule · Top 50 Elite")
    add_day_cards_schedule(prs.slides[-1], schedule_days, footer="3-day VIP program · Vegas 2026")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, TITLE_MAIN, f"Main Event Goals · {EVENT_GOALS_SUBTITLE}")
    add_event_goals_split_view(slide, EVENT_GOALS_MAIN, EVENT_GOALS_DELIVERY)
    set_speaker_notes(slide, EVENT_GOALS)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, TITLE_MAIN, f"Focus Points · {FOCUS_POINTS_SUBTITLE}")
    add_executive_bullets(slide, focus_points)
    set_speaker_notes(slide, focus_points)

    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output


def run(source: Path, pptx_template: Path | None, pptx_out: Path | None) -> dict[str, Path]:
    wb, loaded_from = _open_workbook(source)
    if ROSTER_SHEET not in wb.sheetnames:
        raise ValueError(f"Roster sheet {ROSTER_SHEET!r} not found in {loaded_from}")

    stats = read_top50_stats(wb[ROSTER_SHEET])
    build_schedule_sheet(wb)
    focus_points = build_focus_points_sheet(wb, stats)

    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    export_copy = EXPORT_DIR / "VIP Event - Top 50 Players - schedule-focus.xlsx"
    paths: dict[str, Path] = {}

    wb.save(export_copy)
    paths["export_copy"] = export_copy
    paths["workbook"] = export_copy

    for dest in (loaded_from, DEFAULT_SOURCE, FALLBACK_SOURCE):
        if dest == export_copy or not dest.parent.exists():
            continue
        try:
            shutil.copy2(export_copy, dest)
            paths["workbook"] = dest
            if dest == DEFAULT_SOURCE:
                break
        except PermissionError:
            continue

    out_pptx = pptx_out or EXPORT_DIR / "VIP Event - Vegas 2026 - updated.pptx"
    paths["pptx"] = write_pptx(SCHEDULE_DAYS, focus_points, pptx_template, out_pptx)
    return paths


def main() -> None:
    parser = argparse.ArgumentParser(description="VIP Event Top 50 — Schedule & Focus Points")
    parser.add_argument("--source", type=Path, default=None, help="Top 50 xlsx path")
    parser.add_argument(
        "--pptx",
        type=Path,
        default=None,
        help="Optional VIP Event - Vegas 2026.pptx template (uses blank deck if missing)",
    )
    parser.add_argument(
        "--pptx-out",
        type=Path,
        default=None,
        help="Output pptx path (default: vip_event/exports/...)",
    )
    args = parser.parse_args()

    source = _resolve_source(args.source)
    template = args.pptx if args.pptx else (DEFAULT_PPTX if DEFAULT_PPTX.exists() else None)
    paths = run(source, template, args.pptx_out)

    print(f"Workbook: {paths['workbook']}")
    print(f"  + tabs: {SCHEDULE_SHEET}, {FOCUS_SHEET}")
    print(f"Export copy: {paths['export_copy']}")
    print(f"PowerPoint: {paths['pptx']}")


if __name__ == "__main__":
    main()
