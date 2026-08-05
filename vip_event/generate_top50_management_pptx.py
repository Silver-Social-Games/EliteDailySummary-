"""

Build 2–3 management slides from VIP-Event-Top50-Management.xlsx

and merge into VIP Event - Vegas 2026.pptx (local Desktop VIP Event folder).



Usage:

  python vip_event/generate_top50_management_pptx.py

"""

from __future__ import annotations



import argparse

import re

import shutil

import sys

from dataclasses import dataclass, field

from datetime import date, datetime

from pathlib import Path



from openpyxl import load_workbook

from pptx import Presentation

from pptx.util import Inches, Pt



MODULE_DIR = Path(__file__).resolve().parent

sys.path.insert(0, str(MODULE_DIR))



from pptx_theme import (
    GOLD,
    MUTED,
    TITLE_MAIN,
    add_day_cards_schedule,
    add_event_goals_split_view,
    add_executive_bullets,
    add_kpi_row,
    add_pie_chart,
    add_title_band,
    cohort_two_chart_layout,
    delete_first_slides,
    load_template,
    prepend_slides,
    remove_duplicate_schedule_slides,
    rgb,
    set_speaker_notes,
    style_table_body,
    style_table_header,
)

from event_playbook_content import (
    EVENT_GOALS,
    EVENT_GOALS_MAIN,
    EVENT_GOALS_SECONDARY,
    EVENT_GOALS_SUBTITLE,
    FOCUS_POINTS,
    FOCUS_POINTS_SUBTITLE,
)



FALLBACK_XLSX = MODULE_DIR / "exports" / "VIP Event - Top 50 - Management Brief.xlsx"



VIP_EVENT_DIR = Path(

    r"C:\Users\Owner\OneDrive - Silver Social Games\Desktop\VIP\VIP Event"

)

DEFAULT_SOURCE = VIP_EVENT_DIR / "VIP-Event-Top50-Management.xlsx"

DEFAULT_TEMPLATE = VIP_EVENT_DIR / "VIP Event - Vegas 2026.pptx"



SCHEDULE_DAYS: tuple[tuple[str, ...], ...] = (

    (

        "Pick-up from airport",

        "Gathering",

        "Check-in",

        "Free time",

        "Dinner",

        "Sphere (optional)",

    ),

    (

        "Self breakfast (no pricing)",

        "Gala (3 hours: 13:00-16:00)",

        "Free time",

        "Dinner",

        "Show (Music)",

    ),

    (

        "Brunch + closure",

        "Pick up from hotel",

    ),

)



GEO_PIE_TOP_N = 6



# Detail — Full Data column layout (parity with generate_top50_management_brief.py)

DETAIL_HDR_ROW = 9

DETAIL_DATA_START = 10

DETAIL_DATA_END = 59

DET_COL_RANK = 1
DET_COL_NAME = 3
DET_COL_AGE = 4
DET_COL_STATE = 5
DET_COL_REG = 7
DET_COL_SEN_DAYS = 8
DET_COL_SEN_GROUP = 9
DET_COL_NP_LT = 10
DET_COL_NP_30 = 11
DET_COL_NP_60 = 12
DET_COL_HOLD = 14

DETAIL_KPI_VALUE_ROW = 6
DETAIL_KPI_LABEL_ROW = 7

# Detail row 7 label columns (merged pairs use left column)
DETAIL_KPI_COLS = (1, 3, 5, 7, 8, 10, 12)





def _norm(s) -> str:

    return re.sub(r"\s+", " ", str(s or "").replace("\n", " ").strip())





def _players_text(text: str) -> str:

    """Excel uses invitees in places; presentation uses players."""

    return (

        text.replace("invitees", "players")

        .replace("Invitees", "Players")

        .replace("INVITEES", "PLAYERS")

    )





def _find_sheet(wb, *needles: str):

    for name in wb.sheetnames:

        low = name.lower()

        if all(n.lower() in low for n in needles):

            return wb[name]

    return None





def _cell_str(ws, row: int, col: int) -> str:

    return _norm(ws.cell(row, col).value)





def _fmt_money(val) -> str:

    if val is None:

        return "—"

    try:

        return f"${float(val):,.0f}"

    except (TypeError, ValueError):

        return str(val)





def _fmt_pct(val) -> str:

    if val is None:

        return "—"

    try:

        v = float(val)

        if v == 0:

            return "0.0%"

        return f"{v:.1%}" if abs(v) <= 1 else f"{v:.0f}%"

    except (TypeError, ValueError):

        return str(val)





def _fmt_reg_date(val) -> str:

    if val is None:

        return "—"

    if isinstance(val, datetime):

        return val.strftime("%d-%b-%Y")

    if isinstance(val, date):

        return val.strftime("%d-%b-%Y")

    return str(val)





def _count_value(val) -> float:

    if val is None:

        return 0.0

    try:

        return float(val)

    except (TypeError, ValueError):

        return 0.0





def _format_kpi_value(val, label: str) -> str:

    lbl_u = label.upper()

    if isinstance(val, (int, float)):

        if "%" in label or "HOLD" in lbl_u:

            return _fmt_pct(val)

        if any(k in lbl_u for k in ("NP", "PURCHASE", "NET")):

            return _fmt_money(val)

        if isinstance(val, float) and val == int(val):

            return f"{int(val):,}"

        return f"{val:,}" if isinstance(val, int) else str(val)

    return str(val) if val is not None else "—"





_KPI_LABEL_HINTS = (
    "INVITEE", "PLAYER", "PURCHASE", "NET", "HOLD", "ACTIVE", "MEDIAN",
    "TENURE", "COHORT", "TOTAL", "AVG", "PURCHASED", "NP", "LIFETIME", "DAYS",
)


def _looks_like_kpi_label(lbl: str) -> bool:
    u = lbl.upper()
    if u in ("#", "PLAYER", "COHORT", "STATE", "PLAYERS", "SHARE", "AGENT"):
        return False
    if "RANKS" in u:
        return False
    return any(h in u for h in _KPI_LABEL_HINTS)


def _is_kpi_value(val) -> bool:
    if val is None:
        return False
    if isinstance(val, (int, float)):
        return True
    s = str(val).strip()
    return bool(s.startswith("$") or re.match(r"^[\d]", s) or "%" in s)


def _parse_kpi_card_rows(ws, start_row: int = 7, end_row: int = 14) -> list[list[tuple[str, str]]]:
    """Scan value/label row pairs (e.g. rows 7–8 and 10–11 on slide sheets)."""
    groups: list[list[tuple[str, str]]] = []
    r = start_row
    while r < end_row:
        pairs: list[tuple[int, object, str]] = []
        for c in range(1, 17):
            lbl_raw = ws.cell(r + 1, c).value
            if lbl_raw is None or not str(lbl_raw).strip():
                continue
            lbl = _players_text(_norm(lbl_raw))
            if not _looks_like_kpi_label(lbl):
                continue
            val = ws.cell(r, c).value
            if not _is_kpi_value(val):
                continue
            pairs.append((c, val, lbl))
        if not pairs:
            r += 1
            continue
        cards = [(_format_kpi_value(val, lbl), lbl) for _, val, lbl in sorted(pairs, key=lambda x: x[0])]
        groups.append(cards)
        r += 3
        nxt = _cell_str(ws, r, 1).upper()
        if nxt in ("COHORT", "#") or "TENURE" in nxt or "RANKS" in nxt:
            break
    return groups





def _safe_float(val) -> float | None:
    if val is None:
        return None
    try:
        return float(val)
    except (TypeError, ValueError):
        return None


def _average(vals: list[float]) -> float | None:
    return sum(vals) / len(vals) if vals else None


def _median(vals: list[float]) -> float | None:
    if not vals:
        return None
    s = sorted(vals)
    return s[len(s) // 2]


def _seniority_group(days: int | None) -> str:
    if days is None or days < 0:
        return "—"
    if days == 0:
        return "First Day"
    if days <= 7:
        return "First 2-7d"
    if days <= 30:
        return "First 8-30d"
    if days <= 90:
        return "First 31-90d"
    if days <= 180:
        return "First 91-180"
    if days <= 365:
        return "First 181-365"
    return "+1Y"


def _median_seniority_group(ws) -> str:
    days = []
    for r in range(DETAIL_DATA_START, DETAIL_DATA_END + 1):
        if ws.cell(r, DET_COL_RANK).value is None:
            continue
        d = _safe_float(ws.cell(r, DET_COL_SEN_DAYS).value)
        if d is not None:
            days.append(int(d))
    med = _median([float(d) for d in days]) if days else None
    return _seniority_group(int(med) if med is not None else None)


def _top_state_label(ws) -> str:
    cached = ws.cell(DETAIL_KPI_VALUE_ROW, 8).value
    if cached is not None and str(cached).strip():
        return _players_text(_norm(cached))
    counts: dict[str, int] = {}
    for r in range(DETAIL_DATA_START, DETAIL_DATA_END + 1):
        if ws.cell(r, DET_COL_RANK).value is None:
            continue
        state = _cell_str(ws, r, DET_COL_STATE)
        if state:
            counts[state] = counts.get(state, 0) + 1
    if not counts:
        return "—"
    name = max(counts, key=counts.get)
    n = counts[name]
    total = sum(counts.values()) or 50
    return f"{name} · {n} of {total} ({n / total:.0%})"


def _collect_detail_metrics(ws) -> dict:
    np_lt, np_30, np_60, holds, ages = [], [], [], [], []
    for r in range(DETAIL_DATA_START, DETAIL_DATA_END + 1):
        if ws.cell(r, DET_COL_RANK).value is None:
            continue
        for bucket, col in (
            (np_lt, DET_COL_NP_LT),
            (np_30, DET_COL_NP_30),
            (np_60, DET_COL_NP_60),
            (holds, DET_COL_HOLD),
            (ages, DET_COL_AGE),
        ):
            v = _safe_float(ws.cell(r, col).value)
            if v is not None:
                bucket.append(v)
    return {
        "avg_np_lt": _average(np_lt),
        "avg_np_30": _average(np_30),
        "avg_np_60": _average(np_60),
        "avg_hold": _average(holds),
        "median_age": _median(ages),
        "median_seniority_group": _median_seniority_group(ws),
        "top_state": _top_state_label(ws),
    }


def _detail_kpi_cell_value(ws, col: int, metrics: dict, label: str) -> str:
    """Prefer cached Detail row 6; fall back to computed metrics."""
    raw = ws.cell(DETAIL_KPI_VALUE_ROW, col).value
    lbl_u = label.upper()
    if raw is not None and str(raw).strip():
        if "TOP STATE" in lbl_u:
            return _players_text(_norm(raw))
        if "SENIORITY" in lbl_u or "FTP" in lbl_u or "GROUP" in lbl_u:
            text = _norm(raw)
            if "\n" in text:
                text = text.split("\n", 1)[-1].strip()
            elif re.search(r"\d+\s*d\s", text, re.I):
                text = re.sub(r"^[\d,\s]+d\s*", "", text, flags=re.I).strip()
            return text or metrics["median_seniority_group"]
        return _format_kpi_value(raw, label)

    if "LIFETIME" in lbl_u and "HOLD" not in lbl_u:
        return _fmt_money(metrics["avg_np_lt"])
    if "30" in lbl_u and ("PURCHASE" in lbl_u or "NET" in lbl_u):
        return _fmt_money(metrics["avg_np_30"])
    if "60" in lbl_u and ("PURCHASE" in lbl_u or "NET" in lbl_u):
        return _fmt_money(metrics["avg_np_60"])
    if "HOLD" in lbl_u:
        return _fmt_pct(metrics["avg_hold"])
    if "TOP STATE" in lbl_u:
        return metrics["top_state"]
    if "AGE" in lbl_u:
        age = metrics["median_age"]
        return f"{int(age)}" if age is not None else "—"
    if "SENIORITY" in lbl_u or "FTP" in lbl_u or "GROUP" in lbl_u:
        return metrics["median_seniority_group"]
    return "—"


def build_cohort_average_kpis(ws) -> list[tuple[str, str]]:
    """Single KPI strip — averages only (Detail row 6–7 parity with brief)."""
    metrics = _collect_detail_metrics(ws)
    cards: list[tuple[str, str]] = []
    for col in DETAIL_KPI_COLS:
        label = ws.cell(DETAIL_KPI_LABEL_ROW, col).value
        if label is None or not str(label).strip():
            continue
        lbl = _players_text(_norm(label))
        val = _detail_kpi_cell_value(ws, col, metrics, lbl)
        cards.append((val, lbl))
    return cards


def _is_average_kpi_label(lbl: str) -> bool:
    u = lbl.upper()
    if u.startswith("AVG ") or "AVG " in u:
        return True
    if "MEDIAN" in u:
        return True
    if "TOP STATE" in u:
        return True
    return False


def _filter_average_kpi_rows(
    kpi_rows: list[list[tuple[str, str]]],
) -> list[list[tuple[str, str]]]:
    filtered: list[list[tuple[str, str]]] = []
    for row in kpi_rows:
        cards = [(v, l) for v, l in row if _is_average_kpi_label(l)]
        if cards:
            filtered.append(cards)
    return filtered





def _normalize_seniority_title(title: str) -> str:

    t = title.strip().lstrip("·").strip()

    t = re.sub(r"DAYS\s+FROM\s+FTP", "GROUP MIX", t, flags=re.I)

    t = t.replace("TENURE", "SENIORITY")

    if not t or "FTP" in t.upper():

        return "SENIORITY · GROUP MIX"

    return t





@dataclass

class CohortData:

    subtitle: str = "Top 50 Elite Players · Cohort Overview"

    meta: str = ""

    kpi_rows: list[list[tuple[str, str]]] = field(default_factory=list)

    tenure_title: str = "SENIORITY · GROUP MIX"

    tenure_headers: tuple[str, ...] = ("Cohort", "Players", "Share")

    tenure_rows: list[tuple] = field(default_factory=list)

    geo_title: str = "GEOGRAPHY · STATE MIX"

    geo_headers: tuple[str, ...] = ("State", "Players", "Share")

    geo_rows: list[tuple] = field(default_factory=list)





@dataclass

class RosterRow:

    name: str

    reg_date: str

    seniority: str

    np_lt: str

    np_30d: str

    np_60d: str

    hold_pct: str





@dataclass

class RosterData:

    subtitle: str = "TOP 50 Elite Players"

    meta: str = ""

    kpi_rows: list[list[tuple[str, str]]] = field(default_factory=list)

    left_title: str = "PLAYERS 1 – 25"

    right_title: str = "PLAYERS 26 – 50"

    headers: tuple[str, ...] = (

        "Name", "Reg Date", "Seniority", "NP LT", "NP 30d", "NP 60d", "Hold %",

    )

    left: list[RosterRow] = field(default_factory=list)

    right: list[RosterRow] = field(default_factory=list)





def parse_cohort_sheet(ws) -> CohortData:

    data = CohortData()

    data.subtitle = _players_text(_cell_str(ws, 3, 1)) or data.subtitle

    data.meta = _players_text(_cell_str(ws, 4, 1))

    data.kpi_rows = _parse_kpi_card_rows(ws)



    for r in range(10, 40):

        title = _cell_str(ws, r, 1)

        if title and "TENURE" in title.upper():

            data.tenure_title = _normalize_seniority_title(title)

            break

    for r in range(10, 40):

        title = _cell_str(ws, r, 9)

        if title and "GEOGRAPHY" in title.upper():

            data.geo_title = title.strip()

            break



    hdr_row = None

    for r in range(10, 40):

        if _cell_str(ws, r, 1) == "Cohort" and _cell_str(ws, r, 2) == "Players":

            hdr_row = r

            break

    if hdr_row:

        for r in range(hdr_row + 1, hdr_row + 20):

            cohort = _cell_str(ws, r, 1)

            if not cohort or "AGENT" in cohort.upper():

                break

            n = ws.cell(r, 2).value

            share = ws.cell(r, 3).value

            data.tenure_rows.append((cohort, n, share))



        for r in range(hdr_row + 1, hdr_row + 20):

            state = _cell_str(ws, r, 9)

            if not state:

                continue

            n = ws.cell(r, 10).value

            share = ws.cell(r, 11).value

            data.geo_rows.append((state, n, share))



    return data





def _detail_row_to_roster(ws, row: int) -> RosterRow | None:

    rank = ws.cell(row, DET_COL_RANK).value

    if rank is None:

        return None

    try:

        int(rank)

    except (TypeError, ValueError):

        return None

    name = _cell_str(ws, row, DET_COL_NAME)

    if not name:

        return None

    return RosterRow(

        name=name,

        reg_date=_fmt_reg_date(ws.cell(row, DET_COL_REG).value),

        seniority=_cell_str(ws, row, DET_COL_SEN_GROUP) or "—",

        np_lt=_fmt_money(ws.cell(row, DET_COL_NP_LT).value),

        np_30d=_fmt_money(ws.cell(row, DET_COL_NP_30).value),

        np_60d=_fmt_money(ws.cell(row, DET_COL_NP_60).value),

        hold_pct=_fmt_pct(ws.cell(row, DET_COL_HOLD).value),

    )





def parse_detail_roster(ws) -> RosterData:
    data = RosterData()
    data.subtitle = _players_text(_cell_str(ws, 3, 1)) or data.subtitle
    data.meta = _players_text(_cell_str(ws, 4, 1))

    players: list[RosterRow] = []

    for r in range(DETAIL_DATA_START, DETAIL_DATA_END + 1):

        row = _detail_row_to_roster(ws, r)

        if row:

            players.append(row)

    half = max(1, len(players) // 2)

    data.left = players[:half]

    data.right = players[half:]

    return data





def parse_roster_sheet(ws) -> RosterData:

    """Fallback when Detail tab is unavailable."""

    data = RosterData()

    data.subtitle = _players_text(_cell_str(ws, 3, 1)) or data.subtitle

    data.meta = _players_text(_cell_str(ws, 4, 1))

    data.kpi_rows = _parse_kpi_card_rows(ws)

    return data





def load_management_data(source: Path) -> tuple[CohortData, RosterData]:

    wb = load_workbook(source, data_only=True)

    s1 = _find_sheet(wb, "slide", "1") or _find_sheet(wb, "cohort")

    s2 = _find_sheet(wb, "slide", "2") or _find_sheet(wb, "roster")

    detail = _find_sheet(wb, "detail", "full")

    if s1 is None or s2 is None:

        raise ValueError(f"Expected Slide 1 / Slide 2 tabs in {source}; found {wb.sheetnames}")



    cohort = parse_cohort_sheet(s1)

    if detail is not None:

        roster = parse_detail_roster(detail)

        cohort.kpi_rows = [build_cohort_average_kpis(detail)]

        roster.kpi_rows = []

    else:

        roster = parse_roster_sheet(s2)

        cohort.kpi_rows = _filter_average_kpi_rows(cohort.kpi_rows)

        roster.kpi_rows = []

    return cohort, roster





def _geo_pie_data(rows: list[tuple]) -> tuple[list[str], list[float]]:

    """Top N states by count + Other bucket; values are raw counts for chart labels."""

    sorted_rows = sorted(rows, key=lambda r: _count_value(r[1]), reverse=True)

    top = sorted_rows[:GEO_PIE_TOP_N]

    other_n = sum(_count_value(r[1]) for r in sorted_rows[GEO_PIE_TOP_N:])

    cats = [str(r[0]) for r in top]

    vals = [_count_value(r[1]) for r in top]

    if other_n > 0:

        cats.append("Other")

        vals.append(other_n)

    return cats, vals





def _tenure_pie_data(rows: list[tuple]) -> tuple[list[str], list[float]]:

    cats = [str(r[0]) for r in rows if _count_value(r[1]) > 0]

    vals = [_count_value(r[1]) for r in rows if _count_value(r[1]) > 0]

    return cats, vals





def _resolve_template(path: Path) -> Path:

    candidates = [path]

    if path.parent.exists():

        backups = sorted(path.parent.glob("VIP Event - Vegas 2026 - backup-*.pptx"), reverse=True)

        candidates.extend(backups)

    candidates.append(MODULE_DIR / "exports" / "VIP Event - Vegas 2026 - updated.pptx")

    candidates.append(MODULE_DIR / "exports" / "VIP Event - Vegas 2026 - management.pptx")

    for candidate in candidates:

        if not candidate.exists():

            continue

        try:

            Presentation(str(candidate))

            return candidate

        except Exception:

            continue

    raise FileNotFoundError(f"Cannot open PowerPoint template (close PowerPoint): {path}")





def _open_source(path: Path) -> Path:

    for candidate in (path, FALLBACK_XLSX):

        if not candidate.exists():

            continue

        try:

            load_workbook(candidate, read_only=True).close()

            return candidate

        except PermissionError:

            continue

    raise PermissionError(f"Close Excel and retry: {path}")





def _add_meta_line(slide, text: str, top: float = 1.38) -> None:

    if not text:

        return

    mbox = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12), Inches(0.28))

    mp = mbox.text_frame.paragraphs[0]

    mp.text = text

    mp.font.size = Pt(10)

    mp.font.color.rgb = rgb(MUTED)





def _add_kpi_strips(slide, kpi_rows: list[list[tuple[str, str]]], start_top: float = 1.55) -> float:
    top = start_top
    for row_cards in kpi_rows:
        used = add_kpi_row(slide, top, row_cards)
        if used:
            top += used
    return top


def build_slide_cohort(prs: Presentation, data: CohortData) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_title_band(slide, TITLE_MAIN, data.subtitle)
    _add_meta_line(slide, data.meta)

    kpi_bottom = _add_kpi_strips(slide, data.kpi_rows, 1.55)
    geo_cats, geo_vals = _geo_pie_data(data.geo_rows)
    tenure_cats, tenure_vals = _tenure_pie_data(data.tenure_rows)

    left1, left2, chart_top, chart_w, chart_h = cohort_two_chart_layout(kpi_bottom)
    add_pie_chart(slide, left1, chart_top, chart_w, chart_h, data.geo_title, geo_cats, geo_vals)
    add_pie_chart(
        slide, left2, chart_top, chart_w, chart_h, data.tenure_title, tenure_cats, tenure_vals
    )





def _roster_col_widths(n_cols: int, total_width: float) -> tuple[float, ...]:

    if n_cols == 7:

        return (1.45, 0.82, 0.88, 0.78, 0.72, 0.72, 0.55)

    share = total_width / n_cols

    return tuple(share for _ in range(n_cols))





def _add_roster_block(

    slide,

    left: float,

    top: float,

    width: float,

    title: str,

    headers: tuple[str, ...],

    rows: list[RosterRow],

) -> None:

    bar = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.32))

    bar.fill.solid()

    bar.fill.fore_color.rgb = rgb(GOLD)

    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.05), Inches(width), Inches(0.28))

    tp = tbox.text_frame.paragraphs[0]

    tp.text = title.strip()

    tp.font.size = Pt(11)

    tp.font.bold = True

    tp.font.color.rgb = rgb("000000")



    tbl_top = top + 0.38

    row_h = 0.225

    nrows = 1 + len(rows)

    table = slide.shapes.add_table(

        nrows, len(headers), Inches(left), Inches(tbl_top), Inches(width), Inches(row_h * nrows)

    ).table



    col_widths = _roster_col_widths(len(headers), width)

    for j, w in enumerate(col_widths[: len(headers)]):

        table.columns[j].width = Inches(w)



    for j, h in enumerate(headers):

        style_table_header(table.cell(0, j), h)

    for i, row in enumerate(rows):

        vals = (row.name, row.reg_date, row.seniority, row.np_lt, row.np_30d, row.np_60d, row.hold_pct)

        for j, v in enumerate(vals):

            bold = j == 0

            style_table_body(table.cell(i + 1, j), str(v), bold=bold, size=8, alt_row=i % 2 == 1)





def build_slide_roster(prs: Presentation, data: RosterData) -> None:

    layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(layout)

    add_title_band(slide, TITLE_MAIN, data.subtitle)

    _add_meta_line(slide, data.meta)

    roster_top = 1.68

    _add_roster_block(slide, 0.38, roster_top, 6.18, data.left_title, data.headers, data.left)

    _add_roster_block(slide, 6.78, roster_top, 6.18, data.right_title, data.headers, data.right)





def build_slide_schedule(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_title_band(slide, TITLE_MAIN, "Schedule · Top 50 Elite")
    add_day_cards_schedule(slide, SCHEDULE_DAYS, footer="3-day VIP program · Vegas 2026")


def build_slide_event_goals(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_title_band(slide, TITLE_MAIN, f"Main Event Goals · {EVENT_GOALS_SUBTITLE}")
    add_event_goals_split_view(slide, EVENT_GOALS_MAIN, EVENT_GOALS_SECONDARY)
    set_speaker_notes(slide, EVENT_GOALS)


def build_slide_focus_points(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_title_band(slide, TITLE_MAIN, f"Focus Points · {FOCUS_POINTS_SUBTITLE}")
    add_executive_bullets(slide, FOCUS_POINTS)
    set_speaker_notes(slide, FOCUS_POINTS)


def _build_management_slides(prs: Presentation, cohort: CohortData, roster: RosterData, slide_count: int) -> int:
    build_slide_cohort(prs, cohort)
    build_slide_roster(prs, roster)
    if slide_count >= 3:
        build_slide_schedule(prs)
    if slide_count >= 4:
        build_slide_event_goals(prs)
    if slide_count >= 5:
        build_slide_focus_points(prs)
    return slide_count





def run(

    source: Path,

    template: Path,

    *,

    slide_count: int = 5,

    in_place: bool = True,

    standalone: bool = False,

) -> Path:

    src = _open_source(source)

    cohort, roster = load_management_data(src)

    print(f"Source: {src}")

    print(f"Cohort KPI rows: {sum(len(r) for r in cohort.kpi_rows)} cards")

    print(f"Roster: {len(roster.left)}+{len(roster.right)} players, KPI rows: {sum(len(r) for r in roster.kpi_rows)}")



    if standalone:

        prs = load_template(None)

        _build_management_slides(prs, cohort, roster, slide_count)

        out = MODULE_DIR / "exports" / "VIP Event - Vegas 2026 - management-only.pptx"

        out.parent.mkdir(parents=True, exist_ok=True)

        prs.save(str(out))

        print(f"Standalone: {out} ({len(prs.slides)} slides)")

        return out



    if not template.exists() and not any(

        template.parent.glob("VIP Event - Vegas 2026 - backup-*.pptx")

    ):

        print("Template not found — building standalone deck")

        return run(source, template, slide_count=slide_count, standalone=True)



    template = _resolve_template(template)

    print(f"Using template: {template}")



    backup = template.with_name(

        f"VIP Event - Vegas 2026 - backup-{date.today().isoformat()}.pptx"

    )

    if backup != template:

        try:

            shutil.copy2(template, backup)

            print(f"Backup: {backup}")

        except PermissionError:

            print(f"Warning: could not write backup (file in use): {backup}")



    prs = load_template(str(template))

    orig_count = len(prs.slides)



    delete_first_slides(prs, slide_count)

    before = len(prs.slides)

    _build_management_slides(prs, cohort, roster, slide_count)

    added = len(prs.slides) - before

    if added:

        prepend_slides(prs, added)



    removed = remove_duplicate_schedule_slides(prs, keep_through_index=slide_count - 1)

    if removed:

        print(f"Removed {removed} duplicate Schedule slide(s) from deck")



    out = DEFAULT_TEMPLATE

    try:

        prs.save(str(out))

    except (PermissionError, OSError):

        for fallback in (

            template.parent / "VIP Event - Vegas 2026 - management.pptx",

            MODULE_DIR / "exports" / "VIP Event - Vegas 2026 - management.pptx",

        ):

            try:

                fallback.parent.mkdir(parents=True, exist_ok=True)

                prs.save(str(fallback))

                out = fallback

                print(f"Template locked — wrote: {out}")

                break

            except OSError:

                continue

        else:

            raise

    else:

        print(f"Updated: {out} ({orig_count} -> {len(prs.slides)} slides)")



    export_copy = MODULE_DIR / "exports" / "VIP Event - Vegas 2026 - management.pptx"

    try:

        export_copy.parent.mkdir(parents=True, exist_ok=True)

        shutil.copy2(out, export_copy)

        print(f"Export copy: {export_copy}")

    except OSError:

        pass



    return out





def main() -> None:

    parser = argparse.ArgumentParser(description="VIP Event Top 50 → Vegas 2026 PPT slides")

    parser.add_argument("--source", type=Path, default=DEFAULT_SOURCE)

    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)

    parser.add_argument("--slides", type=int, default=5, choices=(2, 3, 4, 5))

    parser.add_argument("--standalone", action="store_true")

    parser.add_argument("--no-in-place", action="store_true")

    args = parser.parse_args()



    path = run(

        args.source,

        args.template,

        slide_count=args.slides,

        in_place=not args.no_in_place,

        standalone=args.standalone,

    )

    try:

        import os

        os.startfile(path)  # noqa: S606

    except OSError:

        pass





if __name__ == "__main__":

    main()


