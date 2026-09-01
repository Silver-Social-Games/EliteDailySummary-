"""Build the C-level AM Brief and Cursor presentation.

The supplied PowerPoint is used only for its master, theme, colors, logos, and
slide size. All original slide content is removed.
"""

from __future__ import annotations

import argparse
import os
import shutil
import subprocess
import tempfile
from pathlib import Path

from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

MODULE_DIR = Path(__file__).resolve().parent
REPO_ROOT = MODULE_DIR.parent

import sys

sys.path.insert(0, str(REPO_ROOT))

from elite_lib.export_paths import mirror_to_cursor  # noqa: E402
from vip_event.pptx_theme import delete_slide, set_speaker_notes  # noqa: E402


DEFAULT_TEMPLATE = Path(r"C:\Users\Owner\Downloads\VIP Event - Vegas 2026_Final.pptx")
DEFAULT_DASHBOARD = Path(
    r"C:\Users\Owner\OneDrive - Silver Social Games\Desktop\VIP"
    r"\Elite_Cursor\AM Brief\elite_am_brief.html"
)
DEFAULT_BEAM_IMAGE = Path(
    r"C:\Users\Owner\.cursor\projects\c-Users-Owner-Downloads-Elite"
    r"\assets\luxury_beam_clean.png"
)
DEFAULT_PILOT_IMAGE = Path(
    r"C:\Users\Owner\.cursor\projects\c-Users-Owner-Downloads-Elite"
    r"\assets\luxury_pilot_path.png"
)
DEFAULT_IMPACT_ICONS = Path(
    r"C:\Users\Owner\.cursor\projects\c-Users-Owner-Downloads-Elite"
    r"\assets\luxury_impact_icons.png"
)
DEFAULT_CLOSING_ICONS = Path(
    r"C:\Users\Owner\.cursor\projects\c-Users-Owner-Downloads-Elite"
    r"\assets\luxury_qa_thanks_icons.png"
)
DEFAULT_OPENING_ICON = Path(
    r"C:\Users\Owner\.cursor\projects\c-Users-Owner-Downloads-Elite"
    r"\assets\luxury_ai_growth_brain.png"
)
DEFAULT_OUTPUT = MODULE_DIR / "exports" / "AM Brief - From Data to Daily Action.pptx"

GOLD = "F0B429"
GOLD_BRIGHT = "FFCD36"
GOLD_PALE = "FFE09A"
WHITE = "FFFFFF"
OFF_WHITE = "FFF4DC"
MUTED = "9A968C"
MUTED_LIGHT = "C8C4BA"
INK = "1C1C1C"
CARD = "232323"
CARD_ALT = "2C2C2C"
CARD_SOFT = "363636"
BLACK = "000000"

SLIDE_W = 10.0
SLIDE_H = 5.625
CONTENT_BOTTOM = 5.18
FONT = "Calibri"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _remove_all_slides(prs: Presentation) -> None:
    while prs.slides:
        delete_slide(prs, 0)


def _clear_placeholders(slide) -> None:
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)


def _new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _clear_placeholders(slide)
    return slide


def _add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: float = 18,
    color: str = WHITE,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.0,
    font: str = FONT,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return box


def _add_rich_title(slide, gold_text: str, white_text: str = "", *, top: float = 0.35) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(8.9), Inches(0.52))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = gold_text + white_text
    r.font.name = FONT
    r.font.size = Pt(25)
    r.font.bold = True
    r.font.color.rgb = rgb(GOLD)


def _rounded_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: str = CARD,
    line: str = CARD_SOFT,
    radius: float = 0.1,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    if shape.adjustments:
        shape.adjustments[0] = radius
    return shape


def _add_pill(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    *,
    fill: str = CARD_ALT,
    text_color: str = OFF_WHITE,
    line: str = GOLD,
    size: float = 11,
):
    _rounded_card(slide, left, top, width, 0.38, fill=fill, line=line, radius=0.25)
    _add_text(
        slide,
        text,
        left + 0.05,
        top + 0.01,
        width - 0.1,
        0.34,
        size=size,
        color=text_color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def _add_signal_line(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(GOLD)
    line.line.width = Pt(1.6)
    line.line.transparency = 25

    for x, y, d in ((x1, y1, 0.11), ((x1 + x2) / 2, (y1 + y2) / 2, 0.08), (x2, y2, 0.11)):
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x - d / 2),
            Inches(y - d / 2),
            Inches(d),
            Inches(d),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(GOLD_BRIGHT)
        dot.line.fill.background()


def _add_number_badge(slide, number: str, left: float, top: float, diameter: float = 0.34) -> None:
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(diameter),
        Inches(diameter),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(GOLD)
    badge.line.fill.background()
    _add_text(
        slide,
        number,
        left,
        top + 0.01,
        diameter,
        diameter - 0.02,
        size=12,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def _edge_path() -> Path | None:
    candidates = (
        Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
        Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
    )
    return next((path for path in candidates if path.exists()), None)


def capture_dashboard(source_html: Path, output_png: Path) -> Path | None:
    """Capture the manager summary without retaining a copied data payload."""
    edge = _edge_path()
    if not edge or not source_html.is_file():
        return None

    output_png.parent.mkdir(parents=True, exist_ok=True)
    passcode = os.environ.get("ELITE_AM_BRIEF_PASSCODE", "elite")
    injection = f"""
<style>
  html {{ zoom: .82; }}
  .main {{ padding-bottom: 36px !important; }}
</style>
<script>
window.addEventListener("load", function () {{
  setTimeout(function () {{
    var input = document.getElementById("gateInput");
    var submit = document.getElementById("gateSubmit");
    if (input && submit) {{
      input.value = {passcode!r};
      submit.click();
    }}
    window.scrollTo(0, 0);
  }}, 250);
}});
</script>
"""

    with tempfile.TemporaryDirectory(prefix="am-brief-capture-") as temp_raw:
        temp_dir = Path(temp_raw)
        capture_html = temp_dir / "dashboard.html"
        content = source_html.read_text(encoding="utf-8")
        marker = "</body>"
        if marker in content:
            content = content.replace(marker, injection + marker, 1)
        else:
            content += injection
        capture_html.write_text(content, encoding="utf-8")

        raw_png = temp_dir / "dashboard_raw.png"
        profile = temp_dir / "edge-profile"
        command = [
            str(edge),
            "--headless=new",
            "--disable-gpu",
            "--hide-scrollbars",
            "--allow-file-access-from-files",
            f"--user-data-dir={profile}",
            "--window-size=1800,1250",
            "--force-device-scale-factor=1",
            "--virtual-time-budget=4500",
            f"--screenshot={raw_png}",
            capture_html.resolve().as_uri(),
        ]
        result = subprocess.run(command, capture_output=True, text=True, timeout=30, check=False)
        if result.returncode != 0 or not raw_png.is_file():
            return None

        with Image.open(raw_png) as image:
            image = image.convert("RGB")
            crop_bottom = min(image.height, 1010)
            image = image.crop((0, 0, image.width, crop_bottom))
            image = ImageEnhance.Sharpness(image).enhance(1.1)
            image.save(output_png, quality=94)
    return output_png


def _add_image_frame(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    shadow = _rounded_card(
        slide,
        left + 0.06,
        top + 0.07,
        width,
        height,
        fill=BLACK,
        line=BLACK,
        radius=0.04,
    )
    shadow.fill.transparency = 35
    frame = _rounded_card(slide, left, top, width, height, fill=CARD_ALT, line=GOLD, radius=0.04)
    frame.line.width = Pt(1.5)

    with Image.open(image_path) as img:
        iw, ih = img.size
    image_ratio = iw / ih
    box_ratio = width / height
    if image_ratio > box_ratio:
        crop = (1 - box_ratio / image_ratio) / 2
        pic = slide.shapes.add_picture(
            str(image_path), Inches(left + 0.04), Inches(top + 0.04), width=Inches(width - 0.08), height=Inches(height - 0.08)
        )
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1 - image_ratio / box_ratio) / 2
        pic = slide.shapes.add_picture(
            str(image_path), Inches(left + 0.04), Inches(top + 0.04), width=Inches(width - 0.08), height=Inches(height - 0.08)
        )
        pic.crop_top = crop
        pic.crop_bottom = crop


def _add_cover_picture(slide, image_path: Path, left: float, top: float, width: float, height: float):
    """Fill a box with an image while keeping its aspect ratio."""
    with Image.open(image_path) as img:
        iw, ih = img.size
    image_ratio = iw / ih
    box_ratio = width / height
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(left),
        Inches(top),
        width=Inches(width),
        height=Inches(height),
    )
    if image_ratio > box_ratio:
        crop = (1 - box_ratio / image_ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1 - image_ratio / box_ratio) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def _crop_icon_strip(strip: Path, count: int, index: int, output: Path) -> Path | None:
    if not strip.is_file():
        return None
    output.parent.mkdir(parents=True, exist_ok=True)
    with Image.open(strip) as image:
        image = image.convert("RGB")
        cell_left = round(image.width * index / count)
        cell_right = round(image.width * (index + 1) / count)
        cell_width = cell_right - cell_left
        side = min(cell_width, image.height)
        center_y = image.height / 2
        top = max(0, round(center_y - side / 2))
        bottom = min(image.height, top + side)
        cropped = image.crop((cell_left, top, cell_right, bottom)).convert("RGBA")
        pixels = cropped.load()
        for y in range(cropped.height):
            for x in range(cropped.width):
                red, green, blue, _ = pixels[x, y]
                brightness = max(red, green, blue)
                alpha = 0 if brightness < 16 else min(255, max(0, (brightness - 10) * 9))
                pixels[x, y] = (red, green, blue, alpha)
        cropped.save(output)
    return output


def _add_impact_icon(
    slide,
    kind: str,
    center_x: float,
    center_y: float,
    *,
    image_path: Path | None = None,
) -> None:
    """Minimal line icon inside a gold-ring medallion."""
    if image_path and image_path.is_file():
        slide.shapes.add_picture(
            str(image_path),
            Inches(center_x - 0.46),
            Inches(center_y - 0.46),
            width=Inches(0.92),
            height=Inches(0.92),
        )
        return

    outer_d = 0.78
    halo = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(center_x - outer_d / 2),
        Inches(center_y - outer_d / 2),
        Inches(outer_d),
        Inches(outer_d),
    )
    halo.fill.solid()
    halo.fill.fore_color.rgb = rgb(CARD_ALT)
    halo.line.color.rgb = rgb(GOLD)
    halo.line.width = Pt(2.2)

    if kind == "time":
        for x2, y2 in ((center_x, center_y - 0.2), (center_x + 0.17, center_y + 0.06)):
            hand = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(center_x),
                Inches(center_y),
                Inches(x2),
                Inches(y2),
            )
            hand.line.color.rgb = rgb(GOLD)
            hand.line.width = Pt(2.2)
    elif kind == "view":
        eye = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(center_x - 0.25),
            Inches(center_y - 0.13),
            Inches(0.5),
            Inches(0.26),
        )
        eye.fill.background()
        eye.line.color.rgb = rgb(GOLD)
        eye.line.width = Pt(1.8)
        pupil = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(center_x - 0.07),
            Inches(center_y - 0.07),
            Inches(0.14),
            Inches(0.14),
        )
        pupil.fill.solid()
        pupil.fill.fore_color.rgb = rgb(GOLD)
        pupil.line.fill.background()
    else:
        action = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            Inches(center_x - 0.22),
            Inches(center_y - 0.12),
            Inches(0.44),
            Inches(0.24),
        )
        action.fill.solid()
        action.fill.fore_color.rgb = rgb(GOLD)
        action.line.fill.background()


def _build_slide_opportunity(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _add_rich_title(slide, "From Data to Daily Action")
    _add_text(
        slide,
        "Cursor is an AI workspace that turns business questions into tools connected to trusted data.",
        1.0,
        0.86,
        8.0,
        0.38,
        size=12.5,
        color=OFF_WHITE,
        align=PP_ALIGN.CENTER,
    )

    left_x, right_x, card_y, card_w, card_h = 0.78, 5.2, 1.38, 3.72, 2.18
    _rounded_card(slide, left_x, card_y, card_w, card_h, fill=CARD, line=CARD_SOFT)
    _rounded_card(slide, right_x, card_y, card_w, card_h, fill=CARD_ALT, line=GOLD)

    _add_text(slide, "BEFORE", left_x + 0.25, card_y + 0.2, 1.4, 0.28, size=11, color=MUTED, bold=True)
    _add_text(
        slide,
        "Repeated checks\nScattered visibility\nSlow request cycles\nUnclear goals progress",
        left_x + 0.25,
        card_y + 0.57,
        3.1,
        1.38,
        size=16,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _add_text(slide, "NOW", right_x + 0.25, card_y + 0.2, 1.4, 0.28, size=11, color=GOLD, bold=True)
    _add_text(
        slide,
        "Team goals\nDaily protocols\nActionable triggers",
        right_x + 0.25,
        card_y + 0.66,
        3.1,
        1.25,
        size=18,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(4.66),
        Inches(card_y + 0.86),
        Inches(0.38),
        Inches(0.34),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(GOLD)
    arrow.line.fill.background()

    icon_dir = MODULE_DIR / "exports" / "presentation_assets"
    icon_paths = tuple(
        _crop_icon_strip(DEFAULT_IMPACT_ICONS, 3, idx, icon_dir / f"impact_{idx + 1}.png")
        for idx in range(3)
    )
    labels = (("time", "Time Back"), ("view", "Clear Visibility"), ("action", "Faster Action"))
    for idx, (kind, head) in enumerate(labels):
        center_x = 1.72 + idx * 3.28
        _add_impact_icon(slide, kind, center_x, 4.13, image_path=icon_paths[idx])
        _add_text(slide, head, center_x - 1.05, 4.65, 2.1, 0.24, size=12, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

    set_speaker_notes(
        slide,
        [
            (
                "Cursor in one sentence",
                "Cursor is an AI workbench that helps me turn business questions into working tools using trusted company data.",
            ),
            (
                "The before state",
                "The same checks were repeated across different systems. Visibility arrived late, and useful triggers often required a new specialist request.",
            ),
            (
                "The change",
                "I connected the workflow to trusted data and tools, planned the definitions and actions, and created one repeatable AM Brief.",
            ),
            (
                "The impact",
                "I am not claiming an invented number of hours. The clear value is less recurring manual work, shared goal visibility, and faster daily action.",
            ),
        ],
    )


def _build_slide_opening(prs: Presentation) -> None:
    slide = _new_slide(prs)
    icon_dir = MODULE_DIR / "exports" / "presentation_assets"
    opening_icon = _crop_icon_strip(DEFAULT_OPENING_ICON, 1, 0, icon_dir / "opening.png")
    if opening_icon:
        slide.shapes.add_picture(
            str(opening_icon),
            Inches(4.18),
            Inches(1.0),
            width=Inches(1.64),
            height=Inches(1.64),
        )
    _add_text(
        slide,
        "Think Faster. Build Smarter.",
        0.7,
        2.78,
        8.6,
        0.95,
        size=36,
        color=GOLD,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_speaker_notes(
        slide,
        [
            (
                "Open",
                "Today I want to show what becomes possible when an operator can connect Cursor to trusted business data and turn recurring work into a practical daily tool.",
            )
        ],
    )


def _build_slide_dashboard(prs: Presentation, dashboard_image: Path | None) -> None:
    slide = _new_slide(prs)
    image_x, image_y, image_w, image_h = 0.8, 0.34, 8.4, 4.72
    frame = _rounded_card(
        slide,
        image_x - 0.06,
        image_y - 0.06,
        image_w + 0.12,
        image_h + 0.12,
        fill=BLACK,
        line=GOLD,
        radius=0.025,
    )
    frame.line.width = Pt(1.5)
    if dashboard_image and dashboard_image.is_file():
        slide.shapes.add_picture(
            str(dashboard_image),
            Inches(image_x),
            Inches(image_y),
            width=Inches(image_w),
            height=Inches(image_h),
        )
    else:
        _rounded_card(slide, image_x, image_y, image_w, image_h, fill=CARD_ALT, line=GOLD)
        _add_text(
            slide,
            "AM BRIEF\n\nManager view · Team goals · Daily triggers",
            image_x + 0.45,
            image_y + 0.7,
            image_w - 0.9,
            image_h - 1.4,
            size=22,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )

    title_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(image_x),
        Inches(image_y),
        Inches(image_w),
        Inches(0.66),
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = rgb(BLACK)
    title_band.fill.transparency = 12
    title_band.line.fill.background()
    _add_text(
        slide,
        "One Dashboard. One Stop Shop.",
        image_x + 0.2,
        image_y + 0.13,
        image_w - 0.4,
        0.38,
        size=24,
        color=GOLD,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    set_speaker_notes(
        slide,
        [
            (
                "Transition to the live dashboard",
                "This is the proof, not a concept slide. I will open the live AM Brief and show how one view supports a manager and each AM.",
            ),
            (
                "First, goals",
                "Start on the Manager Dashboard. Show Elite Snapshot, Team Goals, and the Goals Leaderboard. Explain that the definitions and pace logic are validated, not generated casually.",
            ),
            (
                "Second, daily triggers",
                "Show the cross-AM trigger tiles, then open one AM view. Highlight WoW purchase gaps, pending redemptions, open tickets, locks, and birthdays.",
            ),
            (
                "Third, action",
                "Show that AIDs connect to Looker and ticket IDs connect to Zendesk. Explain that the brief prioritizes the day and supports action, while sensitive or locked accounts remain protected.",
            ),
            (
                "Return to the deck",
                "Close the live view after two to three minutes and return here. Say: what you just saw is the output; the next slide is what makes it reliable and repeatable.",
            ),
        ],
    )


def _build_slide_method(prs: Presentation) -> None:
    slide = _new_slide(prs)
    _add_rich_title(slide, "Start Small. Prove Value. Scale.")

    card_left, card_top, card_width, card_height = 0.52, 1.22, 8.96, 3.16
    _rounded_card(slide, 0.48, 1.18, 9.04, 3.24, fill=BLACK, line=GOLD, radius=0.025)
    if DEFAULT_BEAM_IMAGE.is_file():
        _add_cover_picture(slide, DEFAULT_BEAM_IMAGE, card_left, card_top, card_width, card_height)
    else:
        _rounded_card(slide, card_left, card_top, card_width, card_height, fill=CARD, line=GOLD)

    beam_y = card_top + card_height / 2
    diameter = 0.96
    stages = (
        ("Connect", "Cursor and trusted data"),
        ("Choose", "One repeated task"),
        ("Plan", "Goal, owner, definitions"),
        ("Build", "Draft, test, refine"),
        ("Improve", "Feedback into the next run"),
        ("Scale", "Repeat across teams"),
    )
    first_x, pitch = 1.30, 1.48
    for index, (label, sub) in enumerate(stages):
        center_x = first_x + index * pitch
        node = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(center_x - diameter / 2),
            Inches(beam_y - diameter / 2),
            Inches(diameter),
            Inches(diameter),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = rgb(CARD)
        node.line.color.rgb = rgb(GOLD)
        node.line.width = Pt(2.2)
        _add_text(
            slide,
            label,
            center_x - 0.66,
            beam_y - 0.18,
            1.32,
            0.36,
            size=14,
            color=OFF_WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        _add_text(
            slide,
            sub,
            center_x - 0.74,
            beam_y + diameter / 2 + 0.12,
            1.48,
            0.56,
            size=9.5,
            color=GOLD,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    ask = _rounded_card(slide, 1.42, 4.6, 7.16, 0.42, fill=GOLD, line=GOLD, radius=0.08)
    ask.line.fill.background()
    _add_text(
        slide,
        "Give Teams the Tools to Move Faster.",
        1.62,
        4.65,
        6.76,
        0.28,
        size=13.5,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    set_speaker_notes(
        slide,
        [
            (
                "One loop, start to scale",
                "Walk the six steps as one loop: connect to trusted data, choose one repeated task, plan it properly, build it, improve it with feedback, then scale it to more teams.",
            ),
            (
                "Always give context",
                "Cursor is strongest when it receives the full business context: the goal, the audience, where the trusted data sits, the definitions, the constraints, and the action the output should drive.",
            ),
            (
                "Use a shared glossary",
                "Explain that a shared glossary gives every KPI, status, and trigger one agreed meaning. Add your own glossary example here.",
            ),
            (
                "Connect the existing ecosystem",
                "The workflow can work across the database, Looker, Zendesk, and Google Sheets. Cursor becomes the workbench that connects the logic across tools already used by the business.",
            ),
            (
                "Validate and protect",
                "Direct access should be permissioned and preferably read-only. Definitions, privacy rules, tests, and human review are part of the product.",
            ),
            (
                "Complement specialist teams",
                "This does not replace Topaz or the data team. It reduces the queue for repeatable requests and leaves specialists more room for high-complexity work and quality control.",
            ),
            (
                "Start small, then the ask",
                "Choose workflows with a clear owner, repeated effort, trusted data, and an observable output. My ask is permission to set up accounts and safe connections, then run one focused pilot per team.",
            ),
        ],
    )


def _build_slide_qa(prs: Presentation) -> None:
    slide = _new_slide(prs)
    icon_dir = MODULE_DIR / "exports" / "presentation_assets"
    qa_icon = _crop_icon_strip(DEFAULT_CLOSING_ICONS, 2, 0, icon_dir / "qa.png")
    if qa_icon:
        slide.shapes.add_picture(
            str(qa_icon),
            Inches(4.18),
            Inches(1.06),
            width=Inches(1.64),
            height=Inches(1.64),
        )
    _add_text(
        slide,
        "Q&A",
        0.7,
        2.82,
        8.6,
        0.95,
        size=44,
        color=GOLD,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_speaker_notes(
        slide,
        [
            (
                "Invite discussion",
                "Open the floor. Useful prompts: Which recurring workflow causes the most friction? Where is trusted data already available? Which team would be a good first pilot?",
            )
        ],
    )


def _build_slide_thanks(prs: Presentation) -> None:
    slide = _new_slide(prs)
    icon_dir = MODULE_DIR / "exports" / "presentation_assets"
    thanks_icon = _crop_icon_strip(DEFAULT_CLOSING_ICONS, 2, 1, icon_dir / "thanks.png")
    if thanks_icon:
        slide.shapes.add_picture(
            str(thanks_icon),
            Inches(4.18),
            Inches(1.06),
            width=Inches(1.64),
            height=Inches(1.64),
        )
    _add_text(
        slide,
        "Thank You",
        0.7,
        2.82,
        8.6,
        0.95,
        size=40,
        color=GOLD,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_speaker_notes(
        slide,
        [
            (
                "Close",
                "Thank the group and restate the ask in one sentence: give selected teams the accounts, safe access, and room to test one useful workflow.",
            )
        ],
    )


def build_deck(template: Path, output: Path, dashboard_image: Path | None) -> Path:
    if not template.is_file():
        raise FileNotFoundError(f"Template not found: {template}")

    prs = Presentation(str(template))
    _remove_all_slides(prs)
    _build_slide_opening(prs)
    _build_slide_opportunity(prs)
    _build_slide_dashboard(prs, dashboard_image)
    _build_slide_method(prs)
    _build_slide_qa(prs)
    _build_slide_thanks(prs)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output


def main() -> None:
    parser = argparse.ArgumentParser(description="Build the AM Brief C-level presentation")
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--dashboard-html", type=Path, default=DEFAULT_DASHBOARD)
    parser.add_argument("--dashboard-image", type=Path)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--skip-capture", action="store_true")
    parser.add_argument("--no-mirror", action="store_true")
    args = parser.parse_args()

    image = args.dashboard_image
    if image is None and not args.skip_capture:
        image = MODULE_DIR / "exports" / "am_brief_manager_summary.png"
        captured = capture_dashboard(args.dashboard_html, image)
        if captured:
            print(f"Dashboard visual: {captured}")
        else:
            print("Dashboard visual unavailable; using the designed fallback panel.")
            image = None

    output = build_deck(args.template, args.output, image)
    print(f"Presentation: {output}")
    print("Slides: 6")
    if not args.no_mirror:
        copied = mirror_to_cursor("am_brief", output)
        if copied:
            print(f"Open: {copied[0]}")


if __name__ == "__main__":
    main()
